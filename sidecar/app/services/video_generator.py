"""Video generator - creates synced MP4 from PowerPoint and narration audio."""

import subprocess
import json
from pathlib import Path
from typing import Dict, List, Optional
import tempfile
import shutil
import os
import io


def check_ffmpeg() -> bool:
    """
    Check if FFmpeg is available.
    
    Returns:
        True if FFmpeg is found, False otherwise
    """
    try:
        result = subprocess.run(
            ['ffmpeg', '-version'],
            capture_output=True,
            text=True,
            timeout=5
        )
        return result.returncode == 0
    except (FileNotFoundError, subprocess.TimeoutExpired):
        # Try common Windows installation paths
        common_paths = [
            r"C:\ffmpeg\bin\ffmpeg.exe",
            r"C:\Program Files\ffmpeg\bin\ffmpeg.exe",
            r"C:\Program Files (x86)\ffmpeg\bin\ffmpeg.exe"
        ]
        for path in common_paths:
            if os.path.exists(path):
                return True
        return False


def get_ffmpeg_path() -> str:
    """
    Get the FFmpeg executable path.
    
    Returns:
        Path to ffmpeg executable
    """
    # Try system PATH first
    try:
        result = subprocess.run(
            ['ffmpeg', '-version'],
            capture_output=True,
            timeout=5
        )
        if result.returncode == 0:
            return 'ffmpeg'
    except:
        pass
    
    # Try common Windows paths
    common_paths = [
        r"C:\ffmpeg\bin\ffmpeg.exe",
        r"C:\Program Files\ffmpeg\bin\ffmpeg.exe",
        r"C:\Program Files (x86)\ffmpeg\bin\ffmpeg.exe"
    ]
    
    for path in common_paths:
        if os.path.exists(path):
            return path
    
    return 'ffmpeg'  # Fallback, will error if not found


def pptx_to_images(pptx_path: Path, output_dir: Path, dpi: int = 150) -> List[Path]:
    """
    Convert PowerPoint slides to PNG images using best available method.
    
    Priority order:
    1. Windows PowerPoint COM (best quality)
    2. LibreOffice (good quality, cross-platform)
    3. Aspose.Slides (commercial, good quality)
    4. python-pptx with PIL rendering (basic, fallback)
    
    Args:
        pptx_path: Path to PPTX file
        output_dir: Directory to save slide images
        dpi: DPI for image export (default: 150 for good quality)
        
    Returns:
        List of paths to generated slide images
    """
    output_dir.mkdir(parents=True, exist_ok=True)
    
    # METHOD 1: Try Windows PowerPoint COM (best quality on Windows)
    if os.name == 'nt':  # Windows only
        try:
            print("🎬 Attempting PowerPoint COM conversion (best quality)...")
            return pptx_to_images_com(pptx_path, output_dir, dpi)
        except Exception as e:
            print(f"⚠️ PowerPoint COM failed: {e}")
    
    # METHOD 2: Try LibreOffice (good quality, cross-platform)
    try:
        print("🎬 Attempting LibreOffice conversion...")
        return pptx_to_images_libreoffice(pptx_path, output_dir, dpi)
    except Exception as e:
        print(f"⚠️ LibreOffice conversion failed: {e}")
    
    # METHOD 3: Try python-pptx with enhanced PIL rendering
    try:
        print("🎬 Attempting python-pptx with PIL rendering...")
        return pptx_to_images_pil(pptx_path, output_dir, dpi)
    except Exception as e:
        print(f"⚠️ PIL rendering failed: {e}")
    
    # FALLBACK: Create placeholder images
    print("⚠️ All conversion methods failed, creating placeholder slides...")
    return create_placeholder_slides(pptx_path, output_dir, dpi)


def pptx_to_images_com(pptx_path: Path, output_dir: Path, dpi: int = 150) -> List[Path]:
    """
    Convert PPTX to images using Windows PowerPoint COM automation.
    This provides the best quality as it uses actual PowerPoint rendering.
    
    Args:
        pptx_path: Path to PPTX file
        output_dir: Directory to save images
        dpi: DPI for export
        
    Returns:
        List of paths to generated images
    """
    try:
        import comtypes.client
        
        # Initialize PowerPoint
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        powerpoint.Visible = 1
        
        # Open presentation
        presentation = powerpoint.Presentations.Open(str(pptx_path.resolve()), WithWindow=False)
        
        slide_images = []
        num_slides = presentation.Slides.Count
        
        print(f"  ✓ Opened presentation with {num_slides} slides")
        
        # Export each slide as PNG
        for idx in range(1, num_slides + 1):
            output_path = output_dir / f"slide_{idx:03d}.png"
            
            # Export slide (format 17 = PNG)
            presentation.Slides[idx].Export(
                str(output_path.resolve()),
                "PNG",
                int(1920),  # Width in pixels (16:9 at 1080p)
                int(1080)   # Height in pixels
            )
            
            slide_images.append(output_path)
            print(f"  ✓ Slide {idx} → {output_path.name}")
        
        # Close presentation and quit PowerPoint
        presentation.Close()
        powerpoint.Quit()
        
        print(f"  ✅ Generated {len(slide_images)} slides using PowerPoint COM")
        return slide_images
        
    except ImportError:
        raise Exception("comtypes not installed (pip install comtypes)")
    except Exception as e:
        # Make sure PowerPoint is closed on error
        try:
            powerpoint.Quit()
        except:
            pass
        raise Exception(f"PowerPoint COM automation failed: {e}")


def pptx_to_images_libreoffice(pptx_path: Path, output_dir: Path, dpi: int = 150) -> List[Path]:
    """
    Convert PPTX to images using LibreOffice command-line tools.
    
    Args:
        pptx_path: Path to PPTX file
        output_dir: Directory to save images
        dpi: DPI for export
        
    Returns:
        List of paths to generated images
    """
    # Find LibreOffice installation
    libreoffice_paths = [
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        "/usr/bin/soffice",
        "/usr/bin/libreoffice",
        "soffice",
        "libreoffice"
    ]
    
    soffice_path = None
    for path in libreoffice_paths:
        try:
            result = subprocess.run(
                [path, '--version'],
                capture_output=True,
                timeout=5
            )
            if result.returncode == 0:
                soffice_path = path
                break
        except:
            continue
    
    if not soffice_path:
        raise Exception("LibreOffice not found")
    
    print(f"  ✓ Found LibreOffice at: {soffice_path}")
    
    # Convert PPTX to PDF first (more reliable)
    pdf_path = output_dir / "presentation.pdf"
    subprocess.run([
        soffice_path,
        '--headless',
        '--convert-to', 'pdf',
        '--outdir', str(output_dir),
        str(pptx_path)
    ], check=True, timeout=60)
    
    if not pdf_path.exists():
        raise Exception("PDF conversion failed")
    
    print("  ✓ Converted to PDF, now converting to images...")
    
    # Convert PDF to images using ImageMagick or pdf2image
    return pdf_to_images(pdf_path, output_dir, dpi)


def pptx_to_images_pil(pptx_path: Path, output_dir: Path, dpi: int = 150) -> List[Path]:
    """
    Convert PPTX to images using python-pptx and PIL (basic rendering).
    This is a fallback method with limited rendering capabilities.
    
    Args:
        pptx_path: Path to PPTX file
        output_dir: Directory to save images
        dpi: DPI for export
        
    Returns:
        List of paths to generated images
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from PIL import Image, ImageDraw, ImageFont
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.dml.color import RGBColor
    import textwrap
    
    prs = Presentation(str(pptx_path))
    slide_images = []
    
    # Calculate image dimensions based on slide size and DPI
    width_inches = prs.slide_width.inches
    height_inches = prs.slide_height.inches
    
    width_px = int(width_inches * dpi)
    height_px = int(height_inches * dpi)
    
    # Ensure dimensions are divisible by 2 (required by H.264 encoder)
    if width_px % 2 != 0:
        width_px += 1
    if height_px % 2 != 0:
        height_px += 1
    
    print(f"  Converting {len(prs.slides)} slides to images ({width_px}x{height_px}px @ {dpi} DPI)")
    
    # Load fonts
    try:
        font_title = ImageFont.truetype("arial.ttf", int(dpi * 0.4))
        font_body = ImageFont.truetype("arial.ttf", int(dpi * 0.2))
        font_small = ImageFont.truetype("arial.ttf", int(dpi * 0.15))
    except:
        font_title = ImageFont.load_default()
        font_body = ImageFont.load_default()
        font_small = ImageFont.load_default()
    
    for idx, slide in enumerate(prs.slides):
        slide_num = idx + 1
        output_path = output_dir / f"slide_{slide_num:03d}.png"
        
        # Get background color
        bg_color = (255, 255, 255)  # Default white
        try:
            if slide.background.fill.type == 1:  # Solid fill
                rgb = slide.background.fill.fore_color.rgb
                bg_color = (rgb[0], rgb[1], rgb[2])
        except:
            pass
        
        # Create canvas
        img = Image.new('RGB', (width_px, height_px), color=bg_color)
        draw = ImageDraw.Draw(img)
        
        # Render shapes
        for shape in slide.shapes:
            try:
                # Get shape position and size
                left = int((shape.left.inches / width_inches) * width_px)
                top = int((shape.top.inches / height_inches) * height_px)
                width = int((shape.width.inches / width_inches) * width_px)
                height = int((shape.height.inches / height_inches) * height_px)
                
                # Render text boxes and placeholders
                if hasattr(shape, "text") and shape.text:
                    text = shape.text
                    
                    # Choose font based on shape type
                    if hasattr(shape, "placeholder_format"):
                        # Title placeholders
                        if shape.placeholder_format.type == 1:  # Title
                            font = font_title
                            text_color = (0, 0, 0)
                        else:
                            font = font_body
                            text_color = (50, 50, 50)
                    else:
                        font = font_body
                        text_color = (50, 50, 50)
                    
                    # Get text color if available
                    try:
                        if shape.text_frame.paragraphs:
                            para = shape.text_frame.paragraphs[0]
                            if para.runs:
                                run = para.runs[0]
                                if run.font.color.type == 1:  # RGB color
                                    rgb = run.font.color.rgb
                                    text_color = (rgb[0], rgb[1], rgb[2])
                    except:
                        pass
                    
                    # Wrap text to fit
                    max_chars = max(1, width // (int(dpi * 0.1)))
                    wrapped_lines = textwrap.wrap(text, width=max_chars)
                    
                    # Draw each line
                    y_offset = top + 10
                    for line in wrapped_lines[:10]:  # Limit to 10 lines
                        draw.text((left + 10, y_offset), line, fill=text_color, font=font)
                        y_offset += int(dpi * 0.25)
                
                # Render images in shapes
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        # Extract image from shape
                        image_stream = shape.image.blob
                        shape_img = Image.open(io.BytesIO(image_stream))
                        
                        # Resize to fit shape dimensions
                        shape_img = shape_img.resize((width, height), Image.LANCZOS)
                        
                        # Paste onto slide
                        img.paste(shape_img, (left, top))
                    except Exception as e:
                        print(f"    ⚠️ Failed to render image in slide {slide_num}: {e}")
            
            except Exception as e:
                print(f"    ⚠️ Failed to render shape in slide {slide_num}: {e}")
                continue
        
        # Save the image
        img.save(str(output_path), 'PNG', quality=95)
        slide_images.append(output_path)
        print(f"  ✓ Slide {slide_num} → {output_path.name}")
    
    print(f"  ✅ Generated {len(slide_images)} slides using PIL rendering")
    return slide_images



def pdf_to_images(pdf_path: Path, output_dir: Path, dpi: int = 150) -> List[Path]:
    """
    Convert PDF to images using FFmpeg or ImageMagick.
    
    Args:
        pdf_path: Path to PDF file
        output_dir: Directory to save images
        dpi: DPI for image export
        
    Returns:
        List of paths to generated images
    """
    try:
        # Try ImageMagick (if available)
        subprocess.run([
            'magick',
            '-density', str(dpi),
            str(pdf_path),
            '-quality', '95',
            str(output_dir / 'slide_%03d.png')
        ], check=True, timeout=120)
        
        # Find generated images
        slide_images = sorted(output_dir.glob('slide_*.png'))
        print(f"  ✓ Generated {len(slide_images)} slide images")
        return slide_images
        
    except:
        print("  ⚠️ ImageMagick not available, trying alternative method...")
        
        # Alternative: Use pdf2image if available
        try:
            from pdf2image import convert_from_path
            
            images = convert_from_path(str(pdf_path), dpi=dpi)
            slide_images = []
            
            for idx, img in enumerate(images):
                output_path = output_dir / f"slide_{idx+1:03d}.png"
                img.save(str(output_path), 'PNG')
                slide_images.append(output_path)
            
            print(f"  ✓ Generated {len(slide_images)} slide images using pdf2image")
            return slide_images
            
        except ImportError:
            raise Exception("Neither ImageMagick nor pdf2image available for PDF conversion")


def create_placeholder_slides(pptx_path: Path, output_dir: Path, dpi: int = 150) -> List[Path]:
    """
    Create placeholder slide images when conversion fails.
    
    Args:
        pptx_path: Path to PPTX file (to count slides)
        output_dir: Directory to save images
        dpi: DPI for images
        
    Returns:
        List of paths to placeholder images
    """
    from pptx import Presentation
    from PIL import Image, ImageDraw, ImageFont
    
    try:
        prs = Presentation(str(pptx_path))
        num_slides = len(prs.slides)
    except:
        num_slides = 10  # Default fallback
    
    slide_images = []
    width_px = int(10 * dpi)  # 10 inches @ dpi
    height_px = int(5.625 * dpi)  # 16:9 aspect ratio
    
    for idx in range(num_slides):
        slide_num = idx + 1
        output_path = output_dir / f"slide_{slide_num:03d}.png"
        
        # Create gradient background
        img = Image.new('RGB', (width_px, height_px), color='#2563eb')
        draw = ImageDraw.Draw(img)
        
        # Add slide number
        try:
            font_large = ImageFont.truetype("arial.ttf", int(dpi * 0.8))
            font_small = ImageFont.truetype("arial.ttf", int(dpi * 0.3))
        except:
            font_large = ImageFont.load_default()
            font_small = ImageFont.load_default()
        
        text = f"Slide {slide_num}"
        
        # Center the text
        bbox = draw.textbbox((0, 0), text, font=font_large)
        text_width = bbox[2] - bbox[0]
        text_height = bbox[3] - bbox[1]
        
        x = (width_px - text_width) // 2
        y = (height_px - text_height) // 2
        
        draw.text((x, y), text, fill='white', font=font_large)
        
        img.save(str(output_path), 'PNG', quality=95)
        slide_images.append(output_path)
    
    print(f"  ✓ Created {len(slide_images)} placeholder slides")
    return slide_images


def generate_srt_subtitles(sentence_timings: List[Dict], output_path: Path) -> Path:
    """
    Generate SRT subtitle file from sentence timings.
    
    Args:
        sentence_timings: List of sentence timing dicts with 'text', 'start', 'end'
        output_path: Path to save .srt file
        
    Returns:
        Path to created SRT file
    """
    def format_srt_time(seconds: float) -> str:
        """Convert seconds to SRT time format: HH:MM:SS,mmm"""
        hours = int(seconds // 3600)
        minutes = int((seconds % 3600) // 60)
        secs = int(seconds % 60)
        millis = int((seconds % 1) * 1000)
        return f"{hours:02d}:{minutes:02d}:{secs:02d},{millis:03d}"
    
    with open(output_path, 'w', encoding='utf-8') as f:
        for idx, timing in enumerate(sentence_timings, 1):
            # SRT format:
            # 1
            # 00:00:00,000 --> 00:00:05,000
            # Subtitle text
            # (blank line)
            
            start_time = format_srt_time(timing['start'])
            end_time = format_srt_time(timing['end'])
            text = timing['text'].strip()
            
            f.write(f"{idx}\n")
            f.write(f"{start_time} --> {end_time}\n")
            f.write(f"{text}\n")
            f.write("\n")
    
    print(f"  ✓ Generated subtitles: {output_path.name} ({len(sentence_timings)} sentences)")
    return output_path


def create_video_from_slides(
    slide_images: List[Path],
    audio_path: Path,
    slide_timings: Dict,
    output_path: Path,
    sentence_timings: Optional[List[Dict]] = None,
    fps: int = 30,
    video_codec: str = 'libx264',
    audio_codec: str = 'aac',
    crf: int = 23,
    add_subtitles: bool = False
) -> Path:
    """
    Create MP4 video from slide images and audio with perfect synchronization.
    
    Args:
        slide_images: List of paths to slide PNG images
        audio_path: Path to narration MP3 file
        slide_timings: Dict with slide timing data (from map_timings_to_slides)
        output_path: Path to output MP4 file
        sentence_timings: Optional list of sentence timing dicts for subtitles
        fps: Frames per second (default: 30)
        video_codec: Video codec (default: libx264 for H.264)
        audio_codec: Audio codec (default: aac)
        crf: Constant Rate Factor for quality (18-28, lower=better, default: 23)
        add_subtitles: Whether to burn subtitles into video (default: False)
        
    Returns:
        Path to created video file
    """
    if not slide_images:
        raise ValueError("No slide images provided")
    
    if not audio_path.exists():
        raise FileNotFoundError(f"Audio file not found: {audio_path}")
    
    print(f"\n{'='*60}")
    print(f"CREATING SYNCED VIDEO")
    print(f"{'='*60}")
    print(f"Slides: {len(slide_images)}")
    print(f"Audio: {audio_path.name}")
    print(f"Output: {output_path.name}")
    print(f"FPS: {fps}, Codec: {video_codec}, CRF: {crf}")
    
    ffmpeg_path = get_ffmpeg_path()
    temp_dir = Path(tempfile.mkdtemp(prefix='lectra_video_'))
    
    try:
        # Step 1: Create a concat file for FFmpeg to specify slide durations
        concat_file = temp_dir / "concat.txt"
        
        with open(concat_file, 'w', encoding='utf-8') as f:
            for slide_timing in slide_timings['slides']:
                slide_num = slide_timing['slide_number']
                duration = slide_timing['duration']
                
                # Find corresponding image
                if slide_num <= len(slide_images):
                    img_path = slide_images[slide_num - 1]
                    
                    # FFmpeg concat demuxer format
                    # file 'path/to/image.png'
                    # duration seconds
                    f.write(f"file '{img_path.absolute()}'\n")
                    f.write(f"duration {duration}\n")
            
            # Add last image again (FFmpeg concat requirement)
            if slide_images:
                f.write(f"file '{slide_images[-1].absolute()}'\n")
        
        print(f"\n📝 Created concat file: {concat_file}")
        
        # Step 2: Get audio duration
        probe_cmd = [
            ffmpeg_path,
            '-i', str(audio_path),
            '-hide_banner'
        ]
        
        probe_result = subprocess.run(
            probe_cmd,
            capture_output=True,
            text=True
        )
        
        # Parse duration from FFmpeg output
        audio_duration = slide_timings.get('total_duration', 60.0)  # fallback
        
        print(f"\n🎵 Audio duration: {audio_duration:.2f} seconds")
        
        # Step 3: Generate subtitles if requested
        subtitle_file = None
        if add_subtitles and sentence_timings:
            subtitle_file = temp_dir / "subtitles.srt"
            generate_srt_subtitles(sentence_timings, subtitle_file)
            print(f"  ✓ Subtitles enabled")
        else:
            print(f"  ℹ Subtitles disabled")
        
        # Step 4: Create video from images with timed slides
        print(f"\n🎬 Generating video with synced slides...")
        
        video_cmd = [
            ffmpeg_path,
            '-f', 'concat',
            '-safe', '0',
            '-i', str(concat_file),
            '-i', str(audio_path),
        ]
        
        # Add subtitle filter if enabled
        if subtitle_file:
            # Escape Windows paths for FFmpeg
            subtitle_path_escaped = str(subtitle_file.absolute()).replace('\\', '/').replace(':', '\\:')
            video_cmd.extend([
                '-vf', f"subtitles='{subtitle_path_escaped}':charenc=UTF-8:force_style='FontName=Arial,FontSize=24,PrimaryColour=&HFFFFFF&,OutlineColour=&H000000&,BorderStyle=3,Outline=2,Shadow=1,MarginV=30'",
            ])
        
        video_cmd.extend([
            '-c:v', video_codec,
            '-r', str(fps),  # Explicit framerate for concat demuxer
            '-pix_fmt', 'yuv420p',  # Compatibility with most players
            '-crf', str(crf),
            '-preset', 'medium',  # Balance between speed and compression
            '-c:a', audio_codec,
            '-b:a', '192k',  # Audio bitrate
            '-shortest',  # End when shortest stream ends
            '-movflags', '+faststart',  # Enable streaming
            '-y',  # Overwrite output
            str(output_path)
        ])
        
        print(f"\n🔧 FFmpeg command:")
        print(f"   {' '.join(video_cmd)}")
        
        # Run FFmpeg
        result = subprocess.run(
            video_cmd,
            capture_output=True,
            text=True,
            timeout=300  # 5 minute timeout
        )
        
        if result.returncode != 0:
            print(f"\n❌ FFmpeg error:")
            print(result.stderr)
            raise RuntimeError(f"FFmpeg failed with code {result.returncode}")
        
        # Verify output
        if not output_path.exists():
            raise RuntimeError("Video file was not created")
        
        file_size_mb = output_path.stat().st_size / (1024 * 1024)
        print(f"\n✅ Video created successfully!")
        print(f"   Path: {output_path}")
        print(f"   Size: {file_size_mb:.2f} MB")
        print(f"   Duration: ~{audio_duration:.1f} seconds")
        
        return output_path
        
    finally:
        # Clean up temp directory
        try:
            shutil.rmtree(temp_dir)
        except:
            pass


def generate_presentation_video(
    pptx_path: Path,
    audio_path: Path,
    slide_timings_path: Path,
    output_dir: Path,
    output_name: str = "presentation_video.mp4",
    dpi: int = 150,
    fps: int = 30
) -> Path:
    """
    Main function to generate a synced video from PPTX and audio.
    
    Args:
        pptx_path: Path to PowerPoint file
        audio_path: Path to narration MP3
        slide_timings_path: Path to slide_timings.json
        output_dir: Directory to save output video
        output_name: Name of output video file
        dpi: DPI for slide images (higher = better quality, default: 150)
        fps: Frames per second (default: 30)
        
    Returns:
        Path to created video file
    """
    # Check FFmpeg availability
    if not check_ffmpeg():
        raise RuntimeError(
            "FFmpeg not found. Please install FFmpeg and add it to PATH, "
            "or install it to C:\\ffmpeg\\bin\\ffmpeg.exe"
        )
    
    # Load slide timings
    with open(slide_timings_path, 'r', encoding='utf-8') as f:
        slide_timings = json.load(f)
    
    print(f"\n{'='*60}")
    print(f"PRESENTATION VIDEO GENERATION PIPELINE")
    print(f"{'='*60}")
    print(f"Input PPTX: {pptx_path.name}")
    print(f"Audio: {audio_path.name}")
    print(f"Slides: {len(slide_timings.get('slides', []))}")
    print(f"Total Duration: {slide_timings.get('total_duration', 0):.1f}s")
    print(f"Output: {output_name}")
    
    # Step 1: Convert PPTX to images
    images_dir = output_dir / "slide_images"
    images_dir.mkdir(parents=True, exist_ok=True)
    
    print(f"\n📸 Step 1: Converting PPTX to images...")
    slide_images = pptx_to_images(pptx_path, images_dir, dpi=dpi)
    
    if not slide_images:
        raise RuntimeError("Failed to convert PPTX to images")
    
    # Step 2: Create video with synced audio
    output_path = output_dir / output_name
    
    # Extract sentence timings if available
    sentence_timings = slide_timings.get('sentences', [])
    
    print(f"\n🎬 Step 2: Creating synced video...")
    video_path = create_video_from_slides(
        slide_images,
        audio_path,
        slide_timings,
        output_path,
        sentence_timings=sentence_timings,
        fps=fps,
        add_subtitles=False  # Respect default parameter - subtitles disabled by default
    )
    
    print(f"\n{'='*60}")
    print(f"✅ PRESENTATION VIDEO COMPLETED!")
    print(f"{'='*60}")
    print(f"Output: {video_path}")
    print(f"Size: {video_path.stat().st_size / (1024*1024):.2f} MB")
    
    return video_path


if __name__ == "__main__":
    # Test the video generator
    import sys
    
    if len(sys.argv) < 4:
        print("Usage: python video_generator.py <pptx_path> <audio_path> <timings_path> [output_dir]")
        sys.exit(1)
    
    pptx_path = Path(sys.argv[1])
    audio_path = Path(sys.argv[2])
    timings_path = Path(sys.argv[3])
    output_dir = Path(sys.argv[4]) if len(sys.argv) > 4 else Path.cwd()
    
    try:
        video_path = generate_presentation_video(
            pptx_path,
            audio_path,
            timings_path,
            output_dir
        )
        print(f"\n✅ Success! Video saved to: {video_path}")
    except Exception as e:
        print(f"\n❌ Error: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)
