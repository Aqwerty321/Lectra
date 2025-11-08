"""PowerPoint presentation generator with professional formatting, images, and charts."""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT, MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pathlib import Path
from typing import Dict, List, Optional, Tuple
import json
import re
import math

from .slide_layout_engine import SlideLayoutEngine, LayoutConfig, SlideLayout, FontScalingStrategy


# Professional color schemes (Gamma-inspired)
COLOR_SCHEMES = {
    "modern_blue": {
        "primary": RGBColor(41, 98, 255),      # Vibrant blue
        "secondary": RGBColor(99, 102, 241),    # Purple-blue
        "accent": RGBColor(250, 204, 21),       # Yellow
        "text_dark": RGBColor(17, 24, 39),      # Almost black
        "text_light": RGBColor(243, 244, 246),  # Light gray
        "background": RGBColor(255, 255, 255)   # White
    },
    "elegant_purple": {
        "primary": RGBColor(124, 58, 237),
        "secondary": RGBColor(167, 139, 250),
        "accent": RGBColor(236, 72, 153),
        "text_dark": RGBColor(17, 24, 39),
        "text_light": RGBColor(243, 244, 246),
        "background": RGBColor(255, 255, 255)
    },
    "professional_teal": {
        "primary": RGBColor(20, 184, 166),
        "secondary": RGBColor(45, 212, 191),
        "accent": RGBColor(251, 146, 60),
        "text_dark": RGBColor(17, 24, 39),
        "text_light": RGBColor(243, 244, 246),
        "background": RGBColor(255, 255, 255)
    }
}


def apply_modern_theme(prs: Presentation, color_scheme: str = "modern_blue"):
    """Apply modern design theme to presentation."""
    colors = COLOR_SCHEMES.get(color_scheme, COLOR_SCHEMES["modern_blue"])
    
    # Set slide dimensions (16:9 widescreen)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    return colors


def add_gradient_background(slide, colors: Dict):
    """Add subtle gradient background to slide."""
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 45.0
    
    # Two-color gradient
    fill.gradient_stops[0].color.rgb = colors["background"]
    fill.gradient_stops[1].color.rgb = RGBColor(
        min(255, colors["background"].r + 10),
        min(255, colors["background"].g + 10),
        min(255, colors["background"].b + 10)
    )


def create_title_slide_modern(prs: Presentation, title: str, subtitle: str, colors: Dict):
    """Create modern, visually appealing title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add decorative shapes
    left_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0), Inches(0),
        Inches(4), prs.slide_height
    )
    left_shape.fill.solid()
    left_shape.fill.fore_color.rgb = colors["primary"]
    left_shape.line.fill.background()
    
    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(4.5), Inches(1.5),
        Inches(5), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = colors["text_dark"]
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(4.5), Inches(3.2),
        Inches(5), Inches(1)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = colors["secondary"]
    
    # Add decorative accent
    accent_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(4.5), Inches(2.9),
        Inches(2), Inches(0.15)
    )
    accent_shape.fill.solid()
    accent_shape.fill.fore_color.rgb = colors["accent"]
    accent_shape.line.fill.background()
    
    return slide


def split_slide_content(points: List[str], max_bullets: int = 4, max_chars: int = 250, 
                        has_image: bool = False) -> List[List[str]]:
    """
    Split slide content into multiple slides if needed.
    Automatically creates "Part 2", "Part 3" etc. when content is too long.
    More conservative splitting to prevent slides with just 1-2 bullets.
    
    Args:
        points: List of bullet points
        max_bullets: Maximum bullets per slide (default: 4, or 3 with image)
        max_chars: Maximum total characters per slide (default: 250, or 150 with image)
        has_image: Whether the slide has an image (reduces available space)
        
    Returns:
        List of bullet point lists (one per slide)
    """
    # Adjust thresholds based on image presence
    if has_image:
        max_bullets = min(max_bullets, 3)  # Max 3 bullets with image
        max_chars = min(max_chars, 150)     # Max 150 chars with image
    else:
        max_bullets = min(max_bullets, 4)   # Max 4 bullets without image
        max_chars = min(max_chars, 200)     # Max 200 chars without image
    
    # Quick check: if within limits, no split needed
    if len(points) <= max_bullets:
        total_chars = sum(len(p) for p in points)
        if total_chars <= max_chars:
            return [points]  # No splitting needed
    
    # Split into chunks intelligently
    slides = []
    current_slide = []
    current_chars = 0
    
    for point in points:
        point_length = len(point)
        
        # Check if adding this point exceeds limits
        # Split if: too many bullets OR too many chars (with at least 2 bullets already to avoid lonely slides)
        if (len(current_slide) >= max_bullets or 
            (current_chars + point_length > max_chars and len(current_slide) >= 2)):
            # Save current slide and start new one
            slides.append(current_slide)
            current_slide = [point]
            current_chars = point_length
        else:
            current_slide.append(point)
            current_chars += point_length
    
    # Add last slide
    if current_slide:
        # If last slide has only 1 bullet and there's a previous slide, merge them
        if len(current_slide) == 1 and len(slides) > 0:
            slides[-1].extend(current_slide)
        else:
            slides.append(current_slide)
    
    return slides if slides else [points]  # Fallback to original if empty


def fit_text_to_box(text_frame, text: str, min_size: int = 12, max_size: int = 24, is_bullet: bool = True):
    """
    Dynamically adjust font size to fit text within box boundaries.
    Uses auto-sizing with word wrapping and proper line spacing.
    
    Args:
        text_frame: python-pptx TextFrame object
        text: Text content to fit
        min_size: Minimum font size in points (default: 12)
        max_size: Maximum font size in points (default: 24)
        is_bullet: Whether this is bullet text (affects line spacing)
        
    Returns:
        Final font size used
    """
    # Enable word wrapping and auto-sizing to prevent overflow
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    
    # Set text content
    p = text_frame.paragraphs[0]
    p.text = text
    
    # Apply consistent line spacing (1.15 recommended for readability)
    p.line_spacing = 1.15 if is_bullet else 1.2
    
    # Use professional fonts (Calibri or Segoe UI)
    p.font.name = 'Calibri'
    p.font.size = Pt(max_size)
    
    # Auto-size will handle shrinking; we return max_size for reference
    return max_size


def create_content_slide_with_dynamic_layout(
    prs: Presentation,
    title: str,
    points: List[str],
    image_path: Optional[str],
    colors: Dict,
    layout_style: str = "left_content_right_image"
):
    """
    🎨 ENHANCED CONTENT SLIDE CREATOR v2.0 with ALGORITHMIC DYNAMIC LAYOUT ENGINE.
    
    Creates professional, visually balanced slides with:
    ✨ Mathematical text fitting (no overlaps, no huge gaps)
    📐 Automatic pagination when content exceeds available space
    🎯 Perfect vertical spacing with centering
    🎨 Adaptive spacing (sparse/dense content detection)
    📊 Font size auto-adjustment for slight overflows
    💎 Premium bullet styling with gradients
    🖼️ Smart image positioning
    
    Algorithm:
        1. Configure LayoutEngine with slide dimensions
        2. Calculate optimal layouts (may return multiple slides)
        3. For each layout:
           - Create slide with title bar
           - Apply vertical centering (top_padding)
           - Render paragraphs with calculated spacing
           - Add premium bullet graphics
           - Position image on first slide only
    
    Args:
        prs: Presentation object
        title: Slide title
        points: List of bullet points
        image_path: Optional path to image
        colors: Color scheme dict
        layout_style: "left_content_right_image" or "full_width"
        
    Returns:
        List of created slide objects
    """
    # === STEP 1: CONFIGURE LAYOUT ENGINE ===
    # Set up configuration matching PowerPoint dimensions
    # CRITICAL: Adjust slide_width based on image presence to prevent overlapping
    
    # Calculate available content width AND spacing
    if layout_style == "left_content_right_image" and image_path:
        # With image on right: content area is ~4.8" wide
        content_width_inches = 4.8
        content_width_pts = content_width_inches * 72  # = 345.6 pts
        chars_per_line = 45  # REDUCED: Even fewer chars for better line breaks
        para_spacing = 24   # INCREASED: Much more space between bullets (was 20)
        font_size = 13      # Slightly smaller font for tighter space
        line_spacing = 1.25 # INCREASED: More breathing room between lines (was 1.2)
    else:
        # Full width: content area is ~9" wide
        content_width_inches = 9.0
        content_width_pts = content_width_inches * 72  # = 648 pts
        chars_per_line = 80  # Standard for full width
        para_spacing = 16   # Slightly more than standard (was 14)
        font_size = 14      # Standard font size
        line_spacing = 1.15 # Standard line spacing
    
    config = LayoutConfig(
        slide_height=405,  # 5.625" × 72 pts/inch = 405 pts (content area after title)
        slide_width=content_width_pts,   # ADJUSTED: Use actual content width
        margin_top=0,      # Title bar handled separately
        margin_bottom=21.6,  # 0.3" bottom margin
        margin_left=36,    # 0.5" left margin
        margin_right=36,   # 0.5" right margin (note: image takes additional space)
        font_size=font_size,      # ADJUSTED: Based on layout style
        line_spacing=line_spacing,  # ADJUSTED: Based on layout style
        para_spacing=para_spacing,   # ADJUSTED: Much more spacing for image slides
        chars_per_line=chars_per_line,  # Adjusted for layout
        auto_center_vertical=True,
        auto_adjust_font=True,
        prevent_orphans=True,
        use_advanced_metrics=True,  # Enable advanced typography
        optimize_spacing=True,
        sparse_content_multiplier=1.5,
        dense_content_multiplier=0.85,
        font_scaling_strategy=FontScalingStrategy.MODERATE
    )
    
    # === STEP 2: CALCULATE LAYOUTS ===
    engine = SlideLayoutEngine(config)
    layouts = engine.calculate_layouts(points, font_size=font_size)  # Use configured font size
    
    # Print debug log for troubleshooting
    print("\n" + "="*70)
    print(f"🎨 DYNAMIC LAYOUT ENGINE v2.0 - Slide: {title}")
    print("="*70)
    for log_msg in engine.get_debug_log():
        print(log_msg)
    print("="*70 + "\n")
    
    # === STEP 3: CREATE SLIDES FROM LAYOUTS ===
    created_slides = []
    
    for layout_idx, layout in enumerate(layouts):
        slide_layout = prs.slide_layouts[6]  # Blank layout for full control
        slide = prs.slides.add_slide(slide_layout)
        
        # === CALCULATE DYNAMIC TITLE BAR HEIGHT ===
        slide_title = title
        if len(layouts) > 1:
            slide_title = f"{title} │ Part {layout_idx + 1} of {len(layouts)}"
        
        # Estimate title lines (assuming ~50 chars per line at font size 28)
        title_chars = len(slide_title)
        estimated_title_lines = max(1, math.ceil(title_chars / 50))
        
        # Calculate title bar height dynamically
        # Base height: 0.85" for single line
        # Additional: 0.3" per extra line
        title_bar_height = 0.85 + (estimated_title_lines - 1) * 0.3
        title_text_height = title_bar_height - 0.25  # Leave margins
        
        print(f"  📏 Title: '{slide_title[:60]}...' ({title_chars} chars, {estimated_title_lines} lines)")
        print(f"     Title bar height: {title_bar_height:.2f}\"")
        
        # === CREATE PREMIUM TITLE BAR WITH DYNAMIC HEIGHT ===
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            prs.slide_width, Inches(title_bar_height)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = colors["primary"]
        title_bar.line.fill.background()
        
        # Add subtle shadow effect (via layering)
        shadow_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(title_bar_height),
            prs.slide_width, Inches(0.03)
        )
        shadow_bar.fill.solid()
        shadow_bar.fill.fore_color.rgb = RGBColor(0, 0, 0)
        shadow_bar.fill.transparency = 0.85  # 85% transparent = subtle shadow
        shadow_bar.line.fill.background()
        
        # Add title text with proper height
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.15),
            Inches(8.5), Inches(title_text_height)
        )
        title_frame = title_box.text_frame
        title_frame.margin_left = Inches(0.1)
        title_frame.margin_right = Inches(0.1)
        title_frame.margin_top = Inches(0.05)
        title_frame.margin_bottom = Inches(0.05)
        title_frame.word_wrap = True
        title_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        
        # Set title text
        p = title_frame.paragraphs[0]
        p.text = slide_title
        p.font.name = 'Calibri'
        p.font.bold = True
        p.font.color.rgb = colors["text_light"]
        p.line_spacing = 1.1
        
        # Adjust font size based on lines
        if estimated_title_lines == 1:
            p.font.size = Pt(32)
        elif estimated_title_lines == 2:
            p.font.size = Pt(26)
        else:
            p.font.size = Pt(22)
        
        # Add decorative accent line (adjust position based on title height)
        accent_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(title_bar_height - 0.10),
            Inches(1.5), Inches(0.04)
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = colors["accent"]
        accent_line.line.fill.background()
        
        # === LAYOUT CONFIGURATION ===
        # Adjust content start position based on dynamic title bar height
        content_top_base = Inches(title_bar_height + 0.35)  # Start below title bar with margin
        
        # Adjust layout based on image presence
        if layout_style == "left_content_right_image" and image_path and layout_idx == 0:
            # Content on left, image on right
            content_left = Inches(0.5)
            content_width = Inches(4.8)  # CRITICAL: Match LayoutConfig slide_width
            
            image_left = Inches(5.8)
            image_top = content_top_base + Inches(0.2)  # Align with content
            image_width = Inches(3.7)
            image_height = Inches(3.7)
        else:
            # Full width content
            content_left = Inches(0.5)
            content_width = Inches(9)
        
        # === APPLY VERTICAL CENTERING ===
        # top_padding = (H_textbox - H_used) / 2
        content_top = content_top_base + Pt(layout.top_padding)
        
        print(f"  🎨 Rendering slide {layout_idx + 1}: {len(layout.paragraphs)} paragraphs")
        print(f"     Content starts at: {content_top}, Vertical padding: {layout.top_padding:.1f} pts")
        print(f"     Density: {layout.density.value}, Utilization: {layout.utilization_ratio:.1%}")
        print(f"     Optimization: {layout.optimization_applied}")
        
        # === RENDER PARAGRAPHS WITH CALCULATED SPACING ===
        current_y_pts = 0  # Track position in points
        
        for para_idx, (para_text, metric) in enumerate(zip(layout.paragraphs, layout.paragraph_metrics)):
            # Add paragraph spacing before each paragraph (except first)
            if para_idx > 0:
                current_y_pts += layout.para_spacing
                print(f"     Added {layout.para_spacing:.1f} pts spacing")
            
            # Detect sub-bullets (indented)
            is_subpoint = para_text.startswith('  ')
            clean_text = para_text.strip()
            
            # === RENDER PREMIUM BULLET POINT ===
            bullet_size = Inches(0.09) if is_subpoint else Inches(0.13)
            bullet_indent = Inches(0.35) if is_subpoint else Inches(0)
            
            # Create gradient bullet (circle with inner highlight)
            bullet_shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                content_left + bullet_indent,
                content_top + Pt(current_y_pts) + Inches(0.06),
                bullet_size, bullet_size
            )
            bullet_shape.fill.solid()
            
            # Use accent color for main bullets, secondary for sub-bullets
            bullet_color = colors["secondary"] if is_subpoint else colors["accent"]
            bullet_shape.fill.fore_color.rgb = bullet_color
            bullet_shape.line.color.rgb = bullet_color
            bullet_shape.line.width = Pt(0.5)
            
            # Add inner highlight circle for depth effect
            if not is_subpoint:
                highlight_size = Inches(0.05)
                highlight_shape = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    content_left + bullet_indent + Inches(0.02),
                    content_top + Pt(current_y_pts) + Inches(0.08),
                    highlight_size, highlight_size
                )
                highlight_shape.fill.solid()
                highlight_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
                highlight_shape.fill.transparency = 0.5
                highlight_shape.line.fill.background()
            
            # === RENDER TEXT WITH CALCULATED HEIGHT ===
            text_left_offset = Inches(0.28) if not is_subpoint else Inches(0.50)
            text_box_width = content_width - text_left_offset
            
            # Use calculated height from layout engine
            text_box_height = Pt(metric.height_required)
            
            text_box = slide.shapes.add_textbox(
                content_left + text_left_offset,
                content_top + Pt(current_y_pts),
                text_box_width,
                text_box_height
            )
            
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            text_frame.auto_size = MSO_AUTO_SIZE.NONE  # Use calculated height
            text_frame.margin_left = Inches(0)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.02)
            text_frame.margin_bottom = Inches(0.02)
            
            # Apply text with calculated font size
            p = text_frame.paragraphs[0]
            p.text = clean_text
            p.font.name = 'Calibri'
            p.font.size = Pt(metric.font_size)
            p.font.color.rgb = colors["text_dark"]
            p.line_spacing = config.line_spacing
            p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
            
            # Add subtle emphasis for key phrases
            if any(keyword in clean_text.lower() for keyword in ['important', 'key', 'critical', 'note:']):
                p.font.bold = True
            
            print(f"     Para {para_idx + 1}: {metric.estimated_lines} lines, "
                  f"{metric.height_required:.1f} pts @ {metric.font_size:.1f}pt font, "
                  f"{metric.word_count} words")
            
            # Update position for next paragraph
            current_y_pts += metric.height_required
        
        # === ADD IMAGE WITH PREMIUM BORDER (only to first slide if multiple pages) ===
        if image_path and layout_idx == 0 and Path(image_path).exists():
            try:
                # Add white border background
                border_padding = Inches(0.1)
                border_shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    image_left - border_padding,
                    image_top - border_padding,
                    image_width + (border_padding * 2),
                    image_height + (border_padding * 2)
                )
                border_shape.fill.solid()
                border_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
                border_shape.line.color.rgb = colors["secondary"]
                border_shape.line.width = Pt(1.5)
                border_shape.shadow.inherit = False
                
                # Add image
                pic = slide.shapes.add_picture(
                    str(image_path),
                    image_left, image_top,
                    width=image_width, height=image_height
                )
                
                print(f"     ✨ Image added with premium border: {Path(image_path).name}")
            except Exception as e:
                print(f"     ⚠️  Failed to add image: {e}")
        
        # === ADD SLIDE NUMBER FOOTER ===
        footer_text = f"Slide {len(created_slides) + 1}"
        footer_box = slide.shapes.add_textbox(
            Inches(8.5), Inches(7.2),
            Inches(1), Inches(0.3)
        )
        footer_frame = footer_box.text_frame
        footer_p = footer_frame.paragraphs[0]
        footer_p.text = footer_text
        footer_p.font.size = Pt(10)
        footer_p.font.color.rgb = RGBColor(150, 150, 150)
        footer_p.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT
        
        created_slides.append(slide)
    
    print(f"  ✅ Created {len(created_slides)} professional slide(s) for: {title}\n")
    return created_slides


def create_content_slide_with_image(
    prs: Presentation,
    title: str,
    points: List[str],
    image_path: Optional[str],
    colors: Dict,
    layout_style: str = "left_content_right_image"
):
    """
    Create content slide with professional layout and optional image.
    Automatically adapts layout based on image presence and content length.
    Prevents text overlap through dynamic spacing, margins, and font sizing.
    
    Args:
        prs: Presentation object
        title: Slide title text
        points: List of bullet points (strings)
        image_path: Optional path to image file
        colors: Color scheme dictionary
        layout_style: "left_content_right_image" or "full_width_image_top"
    
    Returns:
        Created slide object
    """
    slide_layout = prs.slide_layouts[6]  # Blank layout for full control
    slide = prs.slides.add_slide(slide_layout)
    
    # === TITLE BAR (consistent 0.8" height across all slides) ===
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(0.8)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = colors["primary"]
    title_bar.line.fill.background()
    
    # Add title text with margins and auto-sizing
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.15),
        Inches(8), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.margin_left = Inches(0.1)
    title_frame.margin_right = Inches(0.1)
    
    # Apply title formatting with auto-fit
    fit_text_to_box(title_frame, title, min_size=24, max_size=32, is_bullet=False)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = colors["text_light"]
    title_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    
    # === LAYOUT CONFIGURATION (adapts to image presence) ===
    # Enforce minimum 0.5" margins from slide edges
    margin = Inches(0.5)
    content_top = Inches(1.2)  # Start below title bar
    
    if layout_style == "left_content_right_image" and image_path:
        # Content on left (~48%), image on right (~35%), with spacing
        content_left = margin
        content_width = Inches(4.8)
        
        image_left = Inches(5.8)  # 1" spacing between content and image
        image_top = Inches(1.5)
        image_width = Inches(3.7)
        image_height = Inches(3.7)
    elif layout_style == "full_width_image_top" and image_path:
        # Image at top (centered), content below
        image_left = Inches(1.5)
        image_top = Inches(1)
        image_width = Inches(7)
        image_height = Inches(2.5)
        
        content_left = margin
        content_width = Inches(9)  # Full width minus margins
        content_top = Inches(3.8)  # Positioned below image
    else:
        # Full width content, no image
        content_left = margin
        content_width = Inches(9)  # 10" slide - 1" margins
    
    # === PRE-PROCESS BULLET POINTS ===
    # Clean up formatting and intelligently split long text
    processed_points = []
    for point in points:
        # Remove leading bullet symbols (•, -, *, →, ◆, etc.)
        clean_point = point.strip().lstrip('•-*→➤◆►▪').strip()
        
        # Split extremely long bullets (>180 chars) into main + sub-bullets
        if len(clean_point) > 180:
            # Try splitting by colon (main concept: explanation)
            if ':' in clean_point and clean_point.index(':') < 80:
                parts = clean_point.split(':', 1)
                main_part = parts[0].strip() + ':'
                explanation = parts[1].strip()
                
                processed_points.append(main_part)
                
                # Further split long explanations by sentences
                if len(explanation) > 120 and '. ' in explanation:
                    sentences = [s.strip() for s in explanation.split('. ') if s.strip()]
                    for sent in sentences:
                        if not sent.endswith('.'):
                            sent += '.'
                        # Indent sub-bullets with 2 spaces
                        processed_points.append('  ' + sent)
                else:
                    processed_points.append('  ' + explanation)
            else:
                # No colon - add as-is and rely on word wrapping
                processed_points.append(clean_point)
        else:
            processed_points.append(clean_point)
    
    # === CALCULATE DYNAMIC SPACING ===
    # Count main bullets only (not indented sub-bullets)
    main_bullet_count = sum(1 for p in processed_points if not p.startswith('  '))
    
    # Calculate available vertical space (slide height - title bar - margins)
    # Slide height: 5.625", Title bar: 0.8", Top content: 1.2", Bottom margin: 0.3"
    available_height = 5.625 - 1.2 - 0.3  # = 4.125 inches
    
    # Adjust spacing and font size based on content density and available space
    if main_bullet_count >= 5:
        # Very dense content: very tight spacing
        main_bullet_spacing = 0.60
        sub_bullet_spacing = 0.30
        font_size = 12
    elif main_bullet_count == 4:
        # Dense content: tight spacing
        main_bullet_spacing = 0.75
        sub_bullet_spacing = 0.35
        font_size = 13
    elif main_bullet_count == 3:
        # Moderate content
        main_bullet_spacing = 0.90
        sub_bullet_spacing = 0.42
        font_size = 14
    else:
        # Light content (1-2 bullets): comfortable spacing
        main_bullet_spacing = 1.1
        sub_bullet_spacing = 0.50
        font_size = 15
    
    # === RENDER BULLET POINTS ===
    current_y = 0  # Track vertical position
    
    for idx, point in enumerate(processed_points):
        # Detect indented sub-bullets (start with 2 spaces)
        is_subpoint = point.startswith('  ')
        
        # Add vertical spacing between bullets
        if idx > 0:
            current_y += sub_bullet_spacing if is_subpoint else main_bullet_spacing
        
        # Bullet circle size and horizontal indent
        bullet_size = Inches(0.08) if is_subpoint else Inches(0.12)
        bullet_indent = Inches(0.35) if is_subpoint else Inches(0)
        
        # Add bullet shape (colored circle)
        bullet = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            content_left + bullet_indent,
            content_top + Inches(current_y) + Inches(0.05),  # Slight vertical offset
            bullet_size, bullet_size
        )
        bullet.fill.solid()
        bullet.fill.fore_color.rgb = colors["accent"]
        bullet.line.fill.background()  # No border
        
        # Add text box next to bullet
        text_left_offset = Inches(0.25) if not is_subpoint else Inches(0.50)
        text_box_width = content_width - text_left_offset  # Adjust for bullet + spacing
        
        # Calculate text box height dynamically to prevent overlap
        # Estimate lines needed: char count / (width in chars) / chars per line
        estimated_lines = max(1, len(point.strip()) // 60)
        text_box_height = Inches(min(
            sub_bullet_spacing if is_subpoint else main_bullet_spacing,
            0.2 * estimated_lines  # At least 0.2" per line
        ))
        
        text_box = slide.shapes.add_textbox(
            content_left + text_left_offset,
            content_top + Inches(current_y),
            text_box_width,
            text_box_height
        )
        
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # Critical: enables shrink-to-fit
        
        # Set margins for proper padding
        text_frame.margin_left = Inches(0)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.02)
        text_frame.margin_bottom = Inches(0.02)
        
        # Apply text and formatting
        p = text_frame.paragraphs[0]
        p.text = point.strip()
        p.font.name = 'Calibri'  # Professional font
        p.font.size = Pt(font_size - 1 if is_subpoint else font_size)
        p.font.color.rgb = colors["text_dark"]
        p.line_spacing = 1.1  # Slightly tighter to prevent overlap
        p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    
    # === ADD IMAGE (if provided and valid) ===
    if image_path and Path(image_path).exists():
        try:
            slide.shapes.add_picture(
                str(image_path),
                image_left, image_top,
                width=image_width, height=image_height
            )
        except Exception as e:
            print(f"⚠️ Failed to add image to slide: {e}")
    
    return slide


def detect_chart_data(points: List[str]) -> Optional[Dict]:
    """Detect if bullet points contain data suitable for charts."""
    # Look for patterns like "Item: 45%" or "Category - 123"
    data_pattern = re.compile(r'(.+?)[:|\-]\s*(\d+\.?\d*)\s*(%|percent|units?)?', re.IGNORECASE)
    
    matches = []
    for point in points:
        match = data_pattern.search(point)
        if match:
            label = match.group(1).strip()
            value = float(match.group(2))
            matches.append((label, value))
    
    if len(matches) >= 2:  # Need at least 2 data points
        return {
            "labels": [m[0] for m in matches],
            "values": [m[1] for m in matches]
        }
    
    return None


def create_chart_slide(
    prs: Presentation,
    title: str,
    chart_data: Dict,
    colors: Dict,
    chart_type: str = "column"
):
    """Create slide with data visualization chart."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(0.8)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = colors["primary"]
    title_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.15),
        Inches(8), Inches(0.6)
    )
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = colors["text_light"]
    
    # Add chart
    chart_data_obj = CategoryChartData()
    chart_data_obj.categories = chart_data["labels"]
    chart_data_obj.add_series('Values', chart_data["values"])
    
    x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(3.5)
    
    chart_type_map = {
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "pie": XL_CHART_TYPE.PIE,
        "line": XL_CHART_TYPE.LINE
    }
    
    chart = slide.shapes.add_chart(
        chart_type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED),
        x, y, cx, cy, chart_data_obj
    ).chart
    
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    
    return slide


def create_presentation(
    script: Dict,
    output_path: Path,
    slide_images: Optional[Dict[int, List[Dict]]] = None,
    color_scheme: str = "modern_blue"
) -> Path:
    """
    Create professional PowerPoint presentation with images and charts.
    
    Args:
        script: Complete presentation script with slides
        output_path: Path to save PPTX file
        slide_images: Dict mapping slide index to image info (with local_path)
        color_scheme: Color scheme name
        
    Returns:
        Path to created PPTX file
    """
    prs = Presentation()
    colors = apply_modern_theme(prs, color_scheme)
    
    for idx, slide_data in enumerate(script["slides"]):
        if slide_data["type"] == "title":
            # Create modern title slide
            slide = create_title_slide_modern(
                prs,
                slide_data["title"],
                script["title"],
                colors
            )
            
        else:
            # Get image for this slide if available
            image_path = None
            if slide_images and idx in slide_images and len(slide_images[idx]) > 0:
                image_path = slide_images[idx][0].get("local_path")
            
            # Check if we should create a chart
            chart_data = detect_chart_data(slide_data["points"])
            
            if chart_data and len(slide_data["points"]) == len(chart_data["labels"]):
                # Create chart slide
                slide = create_chart_slide(
                    prs,
                    slide_data["title"],
                    chart_data,
                    colors,
                    chart_type="column"
                )
            else:
                # Split slide content based on image presence
                # With image: max 3 bullets or 150 chars
                # Without image: max 4 bullets or 200 chars
                slide_chunks = split_slide_content(
                    slide_data["points"],
                    max_bullets=4,      # Will be adjusted in function based on has_image
                    max_chars=250,      # Will be adjusted in function based on has_image
                    has_image=(image_path is not None)
                )
                
                # Create slide(s) for this content
                for chunk_idx, chunk_points in enumerate(slide_chunks):
                    # Add suffix to title if split into multiple slides
                    chunk_title = slide_data["title"]
                    if len(slide_chunks) > 1:
                        chunk_title = f"{slide_data['title']} (Part {chunk_idx + 1})"
                    
                    # Only add image to first chunk
                    chunk_image = image_path if chunk_idx == 0 else None
                    
                    # Create content slide(s) using DYNAMIC LAYOUT ENGINE
                    layout_style = "left_content_right_image" if chunk_image else "full_width"
                    created_slides = create_content_slide_with_dynamic_layout(
                        prs,
                        chunk_title,
                        chunk_points,
                        chunk_image,
                        colors,
                        layout_style=layout_style
                    )
                    
                    # Add speaker notes (only to first chunk and first slide)
                    if chunk_idx == 0 and len(created_slides) > 0:
                        notes_slide = created_slides[0].notes_slide
                        notes_slide.notes_text_frame.text = slide_data["speaker_notes"]
        
        # Add speaker notes for title slides
        if slide_data["type"] == "title":
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data.get("speaker_notes", "")
    
    # Save presentation
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    
    return output_path
    """
    Create PowerPoint presentation from script.
    
    Args:
        script: Complete presentation script with slides
        output_path: Path to save PPTX file
        
    Returns:
        Path to created PPTX file
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    for slide_data in script["slides"]:
        if slide_data["type"] == "title":
            # Create title slide
            slide_layout = prs.slide_layouts[0]  # Title slide layout
            slide = prs.slides.add_slide(slide_layout)
            
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = slide_data["title"]
            subtitle.text = script["title"]
            
            # Add speaker notes
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data["speaker_notes"]
            
        else:
            # Create content slide
            slide_layout = prs.slide_layouts[1]  # Title and content layout
            slide = prs.slides.add_slide(slide_layout)
            
            title = slide.shapes.title
            title.text = slide_data["title"]
            
            # Add bullet points
            content = slide.placeholders[1]
            text_frame = content.text_frame
            text_frame.clear()
            
            for i, point in enumerate(slide_data["points"]):
                if i == 0:
                    text_frame.text = point
                else:
                    p = text_frame.add_paragraph()
                    p.text = point
                    p.level = 0
            
            # Add speaker notes
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data["speaker_notes"]
    
    # Save presentation
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    
    return output_path


def map_timings_to_slides(
    script: Dict,
    audio_duration: float,
    output_path: Path
) -> Dict:
    """
    Map audio timing to slides for synchronized playback.
    
    Args:
        script: Complete presentation script
        audio_duration: Total audio duration in seconds
        output_path: Path to save timing JSON
        
    Returns:
        Timing data dict
    """
    num_slides = len(script["slides"])
    
    # Simple even distribution (can be enhanced with actual audio analysis)
    time_per_slide = audio_duration / num_slides
    
    timings = {
        "total_duration": audio_duration,
        "slides": []
    }
    
    for i, slide_data in enumerate(script["slides"]):
        start_time = i * time_per_slide
        end_time = (i + 1) * time_per_slide
        
        timings["slides"].append({
            "slide_number": i + 1,
            "title": slide_data["title"],
            "start_time": round(start_time, 2),
            "end_time": round(end_time, 2),
            "duration": round(time_per_slide, 2)
        })
    
    # Save timing data
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(json.dumps(timings, indent=2), encoding='utf-8')
    
    return timings


def save_slide_timings(timings: Dict, output_path: Path):
    """Save slide timing data to JSON file."""
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(json.dumps(timings, indent=2), encoding='utf-8')
