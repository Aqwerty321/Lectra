"""FastAPI application - main API for LECTRA sidecar."""

from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import Optional
from pathlib import Path
import traceback
import asyncio
import json
import re
from pydub import AudioSegment

from .config import config
from .services.samples import EN_SAMPLE, HI_SAMPLE
from .services import (
    tagging, tag_to_ssml, tts_engine, timing, postgres, 
    slide_generator, pptx_generator, image_fetcher, video_generator
)
from .services.sync_calculator import (
    calculate_slide_timings_from_audio,
    probe_media_duration,
    verify_video_sync,
    map_sentences_to_slides
)
from .services import slide_generator_async, image_fetcher_async, tagging_async
from .utils.profiler import timeit, PerformanceProfiler, Timer

# Initialize FastAPI app
app = FastAPI(title="LECTRA Sidecar", version="1.0.0")

# Add CORS middleware to allow frontend to connect
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Allow all origins for Tauri app
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Initialize job logger
job_logger = postgres.create_logger(config.DATABASE_URL)


class GenerateRequest(BaseModel):
    """Request model for /generate endpoint."""
    project: str
    text: str = ""
    lang: str = "en"  # en or hi
    voice: Optional[str] = None
    use_sample: str = "none"  # none, en, hi
    fallback_rate: str = "-10%"
    fallback_pitch: str = "+0st"


class EstimateRequest(BaseModel):
    """Request model for /estimate endpoint."""
    project: str
    text: str = ""
    lang: str = "en"
    voice: Optional[str] = None


class GenerateQuizRequest(BaseModel):
    """Request model for /generate_quiz endpoint."""
    project: str
    slide_start: int
    slide_end: int
    num_questions: int = 5
    difficulty: str = "medium"


class SubmitAnswerRequest(BaseModel):
    """Request model for /submit_quiz_answer endpoint."""
    project: str
    checkpoint_id: str
    question_index: int
    user_answer: str
    use_sample: str = "none"
    fallback_rate: str = "-10%"


class PresentationRequest(BaseModel):
    """Request model for /generate_presentation endpoint."""
    project: str
    topic: str
    lang: str = "en"  # en or hi
    voice: Optional[str] = None
    fallback_rate: str = "-10%"
    fallback_pitch: str = "+0st"
    generate_video: bool = True  # Whether to generate MP4 video (default: True)


@app.get("/healthz")
async def health_check():
    """Health check endpoint."""
    return {"status": "ok", "service": "lectra-sidecar"}


@app.post("/generate")
async def generate_lecture(request: GenerateRequest):
    """
    Full pipeline: tagging → SSML → audio → timing.
    
    Flow:
    1. Choose text (raw or sample)
    2. Call Ollama to add nuance tags
    3. Convert to SSML (for logging)
    4. Generate audio via EdgeTTS
    5. Estimate timings
    6. Log to database
    7. Return paths and preview
    """
    try:
        # 1. Choose text
        if request.use_sample == "en":
            text = EN_SAMPLE
        elif request.use_sample == "hi":
            text = HI_SAMPLE
        else:
            text = request.text
        
        if not text.strip():
            raise HTTPException(status_code=400, detail="No text provided and no sample selected")
        
        # Determine voice
        voice = request.voice
        if not voice:
            voice = config.DEFAULT_HI_VOICE if request.lang == "hi" else config.DEFAULT_EN_VOICE
        
        # Get project directory
        project_dir = config.get_project_dir(request.project)
        
        # 2. Generate tagged text
        try:
            tagged = tagging.generate_nuanced_text(text)
        except ConnectionError as e:
            raise HTTPException(
                status_code=503,
                detail=f"Ollama not reachable: {str(e)}. Please start Ollama and pull llama3.1:latest"
            )
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Tagging failed: {str(e)}")
        
        # Save tagged text
        tagged_path = project_dir / "tagged.txt"
        tagged_path.write_text(tagged, encoding='utf-8')
        
        # 3. Convert to SSML (for logging)
        ssml = tag_to_ssml.to_ssml(tagged, voice, request.fallback_rate, request.fallback_pitch)
        ssml_path = project_dir / "ssml.xml"
        ssml_path.write_text(ssml, encoding='utf-8')
        
        # 4. Parse to segments
        segments = tag_to_ssml.parse_to_segments(tagged, voice)
        
        # 5. Generate audio
        mp3_path = project_dir / "audio.mp3"
        try:
            await tts_engine.speak_edge_async(segments, voice, mp3_path)
        except Exception as e:
            raise HTTPException(
                status_code=502,
                detail=f"EdgeTTS failed: {str(e)}. Check internet connection or try 'estimate only'"
            )
        
        # 6. Estimate timings
        timings_data = timing.estimate_timings(tagged, request.lang, voice, request.fallback_rate)
        timings_json, vtt_path = timing.save_timings(timings_data, project_dir)
        
        # 7. Log to database
        job_logger.log_job(
            lang=request.lang,
            voice=voice,
            input_chars=len(text),
            est_duration_sec=timings_data['total_duration_sec'],
            out_mp3_path=str(mp3_path),
            status="ok"
        )
        
        # Return response
        return {
            "status": "ok",
            "project_dir": str(project_dir),
            "tagged_preview": tagged[:300] + ("..." if len(tagged) > 300 else ""),
            "ssml_path": str(ssml_path),
            "mp3_path": str(mp3_path),
            "timings_path": str(timings_json),
            "vtt_path": str(vtt_path),
            "duration_sec": timings_data['total_duration_sec'],
            "sentence_count": timings_data['sentence_count']
        }
        
    except HTTPException:
        raise
    except Exception as e:
        # Log error
        job_logger.log_job(
            lang=request.lang,
            voice=voice if voice else "unknown",
            input_chars=len(text) if text else 0,
            est_duration_sec=0,
            out_mp3_path="",
            status=f"error: {str(e)}"
        )
        
        raise HTTPException(status_code=500, detail=f"Generation failed: {str(e)}\n{traceback.format_exc()}")


@app.post("/estimate")
async def estimate_lecture(request: EstimateRequest):
    """
    Timing estimation only (no audio generation).
    Quick way to preview duration without synthesis.
    """
    try:
        # 1. Choose text
        if request.use_sample == "en":
            text = EN_SAMPLE
        elif request.use_sample == "hi":
            text = HI_SAMPLE
        else:
            text = request.text
        
        if not text.strip():
            raise HTTPException(status_code=400, detail="No text provided and no sample selected")
        
        # Determine voice
        voice = request.voice
        if not voice:
            voice = config.DEFAULT_HI_VOICE if request.lang == "hi" else config.DEFAULT_EN_VOICE
        
        # Get project directory
        project_dir = config.get_project_dir(request.project)
        
        # 2. Generate tagged text
        try:
            tagged = tagging.generate_nuanced_text(text)
        except ConnectionError as e:
            raise HTTPException(
                status_code=503,
                detail=f"Ollama not reachable: {str(e)}. Please start Ollama and pull llama3.1:latest"
            )
        
        # Save tagged text
        tagged_path = project_dir / "tagged.txt"
        tagged_path.write_text(tagged, encoding='utf-8')
        
        # 3. Estimate timings
        timings_data = timing.estimate_timings(tagged, request.lang, voice, request.fallback_rate)
        timings_json, vtt_path = timing.save_timings(timings_data, project_dir)
        
        return {
            "status": "ok",
            "project_dir": str(project_dir),
            "tagged_preview": tagged[:300] + ("..." if len(tagged) > 300 else ""),
            "timings_path": str(timings_json),
            "vtt_path": str(vtt_path),
            "duration_sec": timings_data['total_duration_sec'],
            "sentence_count": timings_data['sentence_count']
        }
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Estimation failed: {str(e)}\n{traceback.format_exc()}")


@app.post("/generate_presentation")
async def generate_presentation(request: PresentationRequest):
    """
    Full presentation generation pipeline: topic → outline → slides → PPTX → narration → timing.
    
    OPTIMIZED with STREAMING parallel processing:
    - Streaming slide generation (yield slides one-by-one)
    - Per-slide parallel spawning: TTS + image fetch start immediately
    - Overlapping operations: TTS runs concurrently with script generation
    - Incremental assembly: Build assets as they complete
    
    This approach overlaps the TTS bottleneck (54s) with other operations,
    reducing total time from ~115s to target <60s.
    
    Flow:
    1. Generate presentation outline from topic using Ollama
    2. FOR EACH SLIDE (streaming):
       - Generate slide script
       - IMMEDIATELY spawn: TTS task + image fetch task
       - Accumulate completed (script, images, audio) tuples
    3. Build PPTX from accumulated slides
    4. Merge audio chunks into final narration
    5. Calculate timings and generate video
    """
    try:
        # Reset profiler for this request
        PerformanceProfiler.reset()
        
        with Timer("Total Pipeline"):
            if not request.topic.strip():
                raise HTTPException(status_code=400, detail="No topic provided")
            
            # Determine voice
            voice = request.voice
            if not voice:
                voice = config.DEFAULT_HI_VOICE if request.lang == "hi" else config.DEFAULT_EN_VOICE
            
            # Get project directory
            project_dir = config.get_project_dir(request.project)
            
            # STEP 1: Generate outline (fast, sequential)
            with Timer("Step 1: Generate Outline"):
                try:
                    outline = await slide_generator_async.generate_outline_async(request.topic, lang=request.lang)
                    slide_generator_async.save_outline(outline, project_dir / "outline.json")
                except ConnectionError as e:
                    raise HTTPException(
                        status_code=503,
                        detail=f"Ollama not reachable: {str(e)}. Please start Ollama and pull llama3.1:latest"
                    )
                except Exception as e:
                    raise HTTPException(status_code=500, detail=f"Outline generation failed: {str(e)}")
            
            # STEP 2: STREAMING PIPELINE - Process slides as they're generated
            print("\n🚀 Starting STREAMING pipeline (per-slide parallel execution)...")
            
            with Timer("Step 2: Streaming Slide Processing"):
                # Storage for accumulated results
                slides_data = []  # List of slide dicts (script + metadata)
                audio_tasks = {}  # slide_index -> audio task
                image_tasks = {}  # slide_index -> image task
                audio_paths = {}  # slide_index -> audio file path
                slide_images = {}  # slide_index -> list of image paths
                
                # Audio storage directory
                audio_dir = project_dir / "audio_chunks"
                audio_dir.mkdir(exist_ok=True)
                
                # Semaphores for concurrency control
                max_concurrent_images = 5  # Limit concurrent image downloads
                max_concurrent_tts = 3     # Limit concurrent TTS calls
                image_semaphore = asyncio.Semaphore(max_concurrent_images)
                tts_semaphore = asyncio.Semaphore(max_concurrent_tts)
                
                # Wrapper for semaphore-controlled image fetch
                async def fetch_images_with_limit(slide, project_dir):
                    async with image_semaphore:
                        return await image_fetcher_async.fetch_images_for_slide_standalone(slide, project_dir)
                
                # Wrapper for semaphore-controlled TTS
                async def synthesize_with_limit(narration_text, voice, output_mp3):
                    async with tts_semaphore:
                        return await tts_engine.synthesize_slide_narration(narration_text, voice, output_mp3)
                
                # Stream slides and spawn immediate tasks
                async for slide in slide_generator_async.generate_slides_streaming_async(outline, lang=request.lang):
                    slide_idx = slide["slide_index"]
                    slides_data.append(slide)
                    
                    print(f"[✔ Generated] Slide {slide_idx + 1}: {slide['title']}")
                    
                    # SPAWN IMAGE FETCH TASK (if content slide)
                    if slide["type"] == "content" and slide.get("points"):
                        image_task = asyncio.create_task(
                            fetch_images_with_limit(slide, project_dir)
                        )
                        image_tasks[slide_idx] = image_task
                        print(f"  [⚡ Spawned] Image fetch for slide {slide_idx + 1}")
                    
                    # SPAWN TTS TASK (for all slides with speaker notes)
                    if slide.get("speaker_notes"):
                        audio_path = audio_dir / f"slide_{slide_idx:03d}.mp3"
                        tts_task = asyncio.create_task(
                            synthesize_with_limit(
                                slide["speaker_notes"],
                                voice,
                                audio_path
                            )
                        )
                        audio_tasks[slide_idx] = tts_task
                        audio_paths[slide_idx] = audio_path
                        print(f"  [⚡ Spawned] TTS for slide {slide_idx + 1}")
                
                # WAIT for all image tasks to complete
                print("\n⏳ Waiting for all image fetch tasks...")
                for slide_idx, task in image_tasks.items():
                    try:
                        images = await task
                        slide_images[slide_idx] = images
                        print(f"  [✔ Done] Images for slide {slide_idx + 1}: {len(images)} fetched")
                    except Exception as e:
                        print(f"  [⚠️ Failed] Images for slide {slide_idx + 1}: {e}")
                        slide_images[slide_idx] = []
                
                # WAIT for all TTS tasks to complete
                print("\n⏳ Waiting for all TTS tasks...")
                for slide_idx, task in audio_tasks.items():
                    try:
                        await task
                        print(f"  [✔ Done] TTS for slide {slide_idx + 1}")
                    except Exception as e:
                        print(f"  [⚠️ Failed] TTS for slide {slide_idx + 1}: {e}")
            
            # STEP 3: Build complete script structure
            with Timer("Step 3: Assemble Script"):
                script = {
                    "title": outline["title"],
                    "slides": slides_data
                }
                slide_generator_async.save_script(script, project_dir / "script.json")
                
                # Save image metadata
                if slide_images:
                    image_fetcher_async.save_image_metadata(
                        slide_images,
                        project_dir / "images_metadata.json"
                    )
            
            # STEP 4: Create PPTX with all slides and images
            with Timer("Step 4: Create PPTX"):
                try:
                    pptx_path = await asyncio.to_thread(
                        pptx_generator.create_presentation,
                        script,
                        project_dir / "presentation.pptx",
                        slide_images=slide_images,
                        color_scheme="modern_blue"
                    )
                except Exception as e:
                    raise HTTPException(status_code=500, detail=f"PPTX creation failed: {str(e)}")
            
            # STEP 5: Merge audio chunks into final narration
            with Timer("Step 5: Merge Audio Chunks"):
                combined_audio = AudioSegment.empty()
                narration_parts = []
                
                # Merge in slide order
                for slide_idx in sorted(audio_paths.keys()):
                    audio_path = audio_paths[slide_idx]
                    
                    if audio_path.exists():
                        try:
                            audio_chunk = AudioSegment.from_mp3(str(audio_path))
                            combined_audio += audio_chunk
                            
                            # Also accumulate narration text for timing
                            slide = slides_data[slide_idx]
                            if slide.get("speaker_notes"):
                                narration_parts.append(slide["speaker_notes"])
                        except Exception as e:
                            print(f"⚠️ Failed to load audio for slide {slide_idx + 1}: {e}")
                
                if len(combined_audio) == 0:
                    raise HTTPException(
                        status_code=500,
                        detail="No audio was generated"
                    )
                
                # Export final narration
                mp3_path = project_dir / "narration.mp3"
                combined_audio.export(str(mp3_path), format="mp3", bitrate="128k")
                
                # Save narration text
                narration_text = " ".join(narration_parts)
                narration_path = project_dir / "narration.txt"
                narration_path.write_text(narration_text, encoding='utf-8')
                
                audio_duration = len(combined_audio) / 1000.0  # Convert ms to seconds
            
            # STEP 6: Calculate ROBUST slide timings from actual audio
            with Timer("Step 6: Calculate Robust Timings"):
                # First, get estimated timings from text
                estimated_timings = timing.estimate_timings(
                    narration_text,
                    lang=request.lang,
                    default_voice=voice
                )
                
                print(f"📊 Total sentences in narration: {len(estimated_timings['sentences'])}")
                print(f"📊 Total slides: {len(slides_data)}")
                
                # Create BETTER slide mapping based on actual narration structure
                # CRITICAL: Title slides should have minimal duration, not steal sentences!
                slide_mapping = []
                accumulated_text = ""
                sentence_idx = 0
                
                for slide_idx, slide in enumerate(slides_data):
                    slide_num = slide_idx + 1
                    slide_type = slide.get("type", "content")
                    speaker_notes = slide.get("speaker_notes", "")
                    
                    # SPECIAL HANDLING for title slides
                    if slide_type == "title":
                        # Title slides should have SHORT duration (3-5 seconds)
                        # Don't assign any sentences - they'll get fixed duration
                        slide_mapping.append({
                            'slide_number': slide_num,
                            'sentence_indices': [],
                            'is_title': True,
                            'fixed_duration': 4.0  # Fixed 4 seconds for title slides
                        })
                        print(f"  Slide {slide_num} (TITLE): Fixed 4.0s duration (no narration)")
                        continue
                    
                    # For content slides, map sentences based on actual narration
                    accumulated_text += " " + speaker_notes
                    
                    # Count sentences in accumulated text so far
                    # Support both English and Hindi punctuation
                    accumulated_sentences = (
                        accumulated_text.count('.') +   # English period
                        accumulated_text.count('?') +   # Question mark
                        accumulated_text.count('!') +   # Exclamation mark
                        accumulated_text.count('।') +   # Hindi purna viram (Devanagari full stop)
                        accumulated_text.count('॥')    # Hindi double danda (end of verse/section)
                    )
                    
                    # Determine how many sentences this slide should get
                    target_sentence_count = accumulated_sentences - sentence_idx
                    
                    indices = []
                    for _ in range(target_sentence_count):
                        if sentence_idx < len(estimated_timings['sentences']):
                            indices.append(sentence_idx)
                            sentence_idx += 1
                        else:
                            break
                    
                    # Only add content slides that have narration
                    if indices:
                        slide_mapping.append({
                            'slide_number': slide_num,
                            'sentence_indices': indices,
                            'is_title': False
                        })
                        print(f"  Slide {slide_num} (CONTENT): {len(indices)} sentences (indices {indices[0]}-{indices[-1]})")
                    else:
                        # Content slide with no narration - give it a brief duration
                        slide_mapping.append({
                            'slide_number': slide_num,
                            'sentence_indices': [],
                            'is_title': False,
                            'fixed_duration': 2.0
                        })
                        print(f"  Slide {slide_num} (CONTENT): Fixed 2.0s duration (no narration)")
                
                # Handle remaining sentences (distribute to last content slide)
                if sentence_idx < len(estimated_timings['sentences']):
                    remaining = list(range(sentence_idx, len(estimated_timings['sentences'])))
                    # Find last content slide
                    for mapping in reversed(slide_mapping):
                        if not mapping.get('is_title', False):
                            mapping['sentence_indices'].extend(remaining)
                            print(f"  ⚠️  Added {len(remaining)} remaining sentences to slide {mapping['slide_number']}")
                            break
                
                # Calculate ROBUST timings from ACTUAL audio file
                slide_timings = calculate_slide_timings_from_audio(
                    audio_path=mp3_path,
                    sentence_timings=estimated_timings['sentences'],
                    slide_mapping=slide_mapping
                )
                
                # Save slide timings
                timings_json_path = project_dir / "slide_timings.json"
                with open(timings_json_path, 'w', encoding='utf-8') as f:
                    json.dump(slide_timings, f, indent=2)
                
                print(f"✅ Robust timings calculated: {slide_timings['actual_duration']:.3f}s actual vs {slide_timings['estimated_duration']:.3f}s estimated")
                print(f"   Scale factor: {slide_timings['scale_factor']:.4f} ({len(slide_timings['slides'])} slides)")
                
                # Log to database
                job_logger.log_job(
                    lang=request.lang,
                    voice=voice,
                    input_chars=len(narration_text),
                    est_duration_sec=slide_timings['actual_duration'],  # Use ACTUAL duration
                    out_mp3_path=str(mp3_path),
                    status="ok"
                )
            
            # STEP 6.5: Generate Interactive Animations (NEW!)
            with Timer("Step 6.5: Generate Animations"):
                try:
                    print("\n🎨 Generating interactive animations for slides...")
                    from .services.ollama_client import generate_animation_steps
                    
                    animations_data = {
                        "slides": []
                    }
                    
                    for slide_idx, slide in enumerate(slides_data):
                        slide_num = slide_idx + 1
                        slide_type = slide.get("type", "content")
                        
                        # Skip title slides for animations (they're static)
                        if slide_type == "title":
                            animations_data["slides"].append({
                                "slide_number": slide_num,
                                "steps": [{
                                    "id": 1,
                                    "text": slide.get("title", ""),
                                    "action": "fadeIn",
                                    "duration": 2.0,
                                    "hint": "Welcome to this lecture!",
                                    "element": "title"
                                }]
                            })
                            continue
                        
                        # Generate animations for content slides
                        slide_content = "\n".join(slide.get("points", []))
                        slide_title = slide.get("title", "")
                        speaker_notes = slide.get("speaker_notes", "")
                        
                        try:
                            animation_steps = generate_animation_steps(
                                slide_content=slide_content,
                                slide_title=slide_title,
                                speaker_notes=speaker_notes
                            )
                            
                            # Add timing information from slide_timings
                            slide_timing = next(
                                (s for s in slide_timings['slides'] if s['slide'] == slide_idx),
                                None
                            )
                            
                            animations_data["slides"].append({
                                "slide_number": slide_num,
                                "start_time": slide_timing['start'] if slide_timing else 0,
                                "end_time": slide_timing['end'] if slide_timing else 0,
                                "steps": animation_steps.get("steps", [])
                            })
                            
                            print(f"  ✓ Generated {len(animation_steps.get('steps', []))} animation steps for slide {slide_num}")
                            
                        except Exception as e:
                            print(f"  ⚠️ Animation generation failed for slide {slide_num}: {e}")
                            # Add fallback simple animation
                            animations_data["slides"].append({
                                "slide_number": slide_num,
                                "steps": [{
                                    "id": 1,
                                    "text": slide_title,
                                    "action": "fadeIn",
                                    "duration": 2.0,
                                    "hint": "Focus on this key concept",
                                    "element": "text"
                                }]
                            })
                    
                    # Save animations JSON
                    animations_path = project_dir / "animations.json"
                    with open(animations_path, 'w', encoding='utf-8') as f:
                        json.dump(animations_data, f, indent=2)
                    
                    print(f"✅ Animations saved: {animations_path}")
                    
                except Exception as e:
                    print(f"⚠️ Animation generation failed (lecture still available): {str(e)}")
                    animations_path = None
            
            # STEP 7: Generate synced video (optional)
            video_path = None
            if request.generate_video:
                with Timer("Step 7: Generate Video"):
                    try:
                        print("\n🎬 Generating synced presentation video...")
                        video_path = video_generator.generate_presentation_video(
                            pptx_path=pptx_path,
                            audio_path=mp3_path,
                            slide_timings_path=project_dir / "slide_timings.json",
                            output_dir=project_dir,
                            output_name="presentation_video.mp4",
                            dpi=150,
                            fps=30
                        )
                        print(f"✅ Video created: {video_path}")
                        
                        # Verify sync
                        is_synced, actual_dur, message = verify_video_sync(
                            video_path,
                            slide_timings['actual_duration'],
                            tolerance=0.5
                        )
                        if not is_synced:
                            print(f"⚠️ Warning: {message}")
                        
                    except Exception as e:
                        print(f"⚠️ Video generation failed (presentation still available): {str(e)}")
            
            # Print performance report
            print(PerformanceProfiler.get_report())
            
            # Return response with image count
            image_count = sum(len(imgs) for imgs in slide_images.values())
            
            response_data = {
                "status": "ok",
                "project_dir": str(project_dir),
                "presentation_title": script["title"],
                "slide_count": len(script["slides"]),
                "image_count": image_count,
                "pptx_path": str(pptx_path),
                "audio_path": str(mp3_path),
                "narration_path": str(narration_path),
                "slide_timings_path": str(project_dir / "slide_timings.json"),
                "animations_path": str(project_dir / "animations.json") if (project_dir / "animations.json").exists() else None,
                "duration_sec": audio_duration,
                "outline": outline,
                "slide_timings": slide_timings,
                "features": ["AI-generated content", "Professional design", "Auto-fetched images", "Charts & diagrams", "Synchronized narration", "Streaming parallel execution", "Interactive animations"]
            }
            
            # Add video path if generated
            if video_path:
                response_data["video_path"] = str(video_path)
                response_data["features"].append("MP4 video with synced slides")
            
            return response_data
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Presentation generation failed: {str(e)}\n{traceback.format_exc()}")


@app.post("/generate_presentation_legacy")
async def generate_presentation_legacy(request: PresentationRequest):
    """
    LEGACY VERSION: Full presentation generation pipeline with batch parallel processing.
    
    Use /generate_presentation for the new streaming version which overlaps TTS with other operations.
    
    Flow:
    1. Generate presentation outline from topic using Ollama
    2. Generate content for all slides IN PARALLEL
    3. Fetch images for all slides IN PARALLEL (while PPTX is being created)
    4. Create PowerPoint file with slides
    5. Generate narration script from speaker notes (async)
    6. Synthesize audio with EdgeTTS
    7. Map audio timing to slides
    8. Generate video (if requested)
    """
    try:
        # Reset profiler for this request
        PerformanceProfiler.reset()
        
        with Timer("Total Pipeline"):
            if not request.topic.strip():
                raise HTTPException(status_code=400, detail="No topic provided")
            
            # Determine voice
            voice = request.voice
            if not voice:
                voice = config.DEFAULT_HI_VOICE if request.lang == "hi" else config.DEFAULT_EN_VOICE
            
            # Get project directory
            project_dir = config.get_project_dir(request.project)
            
            # STEP 1: Generate outline (fast, sequential)
            with Timer("Step 1: Generate Outline"):
                try:
                    outline = await slide_generator_async.generate_outline_async(request.topic, lang=request.lang)
                    slide_generator_async.save_outline(outline, project_dir / "outline.json")
                except ConnectionError as e:
                    raise HTTPException(
                        status_code=503,
                        detail=f"Ollama not reachable: {str(e)}. Please start Ollama and pull llama3.1:latest"
                    )
                except Exception as e:
                    raise HTTPException(status_code=500, detail=f"Outline generation failed: {str(e)}")
            
            # STEP 2: Generate full script (PARALLEL - all slides at once)
            with Timer("Step 2: Generate Script (Parallel)"):
                try:
                    script = await slide_generator_async.generate_full_script_async(outline)
                    slide_generator_async.save_script(script, project_dir / "script.json")
                except Exception as e:
                    raise HTTPException(status_code=500, detail=f"Script generation failed: {str(e)}")
            
            # STEP 3: Fetch images (ASYNC - all at once)
            with Timer("Step 3: Fetch Images (Parallel)"):
                try:
                    slide_images = await image_fetcher_async.fetch_images_for_slides_async(script, project_dir)
                    
                    # Save image metadata
                    if slide_images:
                        image_fetcher_async.save_image_metadata(slide_images, project_dir / "images_metadata.json")
                except Exception as e:
                    print(f"⚠️ Image fetching failed (non-critical): {e}")
                    slide_images = {}
            
            # STEP 4: Create PPTX with images (only once!)
            with Timer("Step 4: Create PPTX"):
                try:
                    pptx_path = await asyncio.to_thread(
                        pptx_generator.create_presentation,
                        script,
                        project_dir / "presentation.pptx",
                        slide_images=slide_images,
                        color_scheme="modern_blue"
                    )
                except Exception as e:
                    raise HTTPException(status_code=500, detail=f"PPTX creation failed: {str(e)}")
            
            # STEP 5: Generate narration from speaker notes
            with Timer("Step 5: Build Narration Text"):
                narration_parts = []
                for slide in script["slides"]:
                    if slide.get("speaker_notes"):
                        narration_parts.append(slide["speaker_notes"])
                
                if not narration_parts:
                    raise HTTPException(
                        status_code=400,
                        detail="No speaker notes found in slides to generate narration"
                    )
                
                narration_text = " ".join(narration_parts)
                
                if len(narration_text.strip()) < 10:
                    raise HTTPException(
                        status_code=400,
                        detail="Narration text too short (need at least 10 characters)"
                    )
                
                # Save narration text
                narration_path = project_dir / "narration.txt"
                narration_path.write_text(narration_text, encoding='utf-8')
            
            # STEP 6: Generate tagged narration with prosody (ASYNC)
            with Timer("Step 6: Tag Narration (Async)"):
                try:
                    tagged_narration = await tagging_async.generate_nuanced_text_async(narration_text)
                    tagged_path = project_dir / "narration_tagged.txt"
                    tagged_path.write_text(tagged_narration, encoding='utf-8')
                except Exception as e:
                    raise HTTPException(status_code=500, detail=f"Narration tagging failed: {str(e)}")
            
            # STEP 7: Parse to segments and generate audio
            with Timer("Step 7: Parse Segments"):
                segments = tag_to_ssml.parse_to_segments(tagged_narration, voice)
                
                if not segments:
                    raise HTTPException(
                        status_code=400,
                        detail="Failed to parse narration into audio segments"
                    )
            
            # STEP 8: Synthesize audio with EdgeTTS
            mp3_path = project_dir / "narration.mp3"
            
            with Timer("Step 8: Synthesize Audio (EdgeTTS)"):
                try:
                    await tts_engine.speak_edge_async(segments, voice, mp3_path)
                except Exception as e:
                    error_msg = str(e)
                    if "No audio was received" in error_msg:
                        raise HTTPException(
                            status_code=502,
                            detail=f"EdgeTTS failed: No audio was received. Please verify that your parameters are correct.. Check internet connection or try a different voice (current: {voice})"
                        )
                    raise HTTPException(
                        status_code=502,
                        detail=f"EdgeTTS failed: {error_msg}. Check internet connection"
                    )
            
            # STEP 9: Estimate audio duration and map to slides
            with Timer("Step 9: Calculate Timings"):
                timings_data = timing.estimate_timings(tagged_narration, request.lang, voice, request.fallback_rate)
                audio_duration = timings_data['total_duration_sec']
                
                # Map timings to slides
                slide_timings = pptx_generator.map_timings_to_slides(
                    script,
                    audio_duration,
                    project_dir / "slide_timings.json"
                )
                
                # Log to database
                job_logger.log_job(
                    lang=request.lang,
                    voice=voice,
                    input_chars=len(narration_text),
                    est_duration_sec=audio_duration,
                    out_mp3_path=str(mp3_path),
                    status="ok"
                )
            
            # STEP 10: Generate synced video (optional, in background)
            video_path = None
            if request.generate_video:
                with Timer("Step 10: Generate Video"):
                    try:
                        print("\n🎬 Generating synced presentation video...")
                        video_path = video_generator.generate_presentation_video(
                            pptx_path=pptx_path,
                            audio_path=mp3_path,
                            slide_timings_path=project_dir / "slide_timings.json",
                            output_dir=project_dir,
                            output_name="presentation_video.mp4",
                            dpi=150,
                            fps=30
                        )
                        print(f"✅ Video created: {video_path}")
                    except Exception as e:
                        print(f"⚠️ Video generation failed (presentation still available): {str(e)}")
            
            # Print performance report
            print(PerformanceProfiler.get_report())
            
            # Return response with image count
            image_count = sum(len(imgs) for imgs in slide_images.values())
            
            response_data = {
                "status": "ok",
                "project_dir": str(project_dir),
                "presentation_title": script["title"],
                "slide_count": len(script["slides"]),
                "image_count": image_count,
                "pptx_path": str(pptx_path),
                "audio_path": str(mp3_path),
                "narration_path": str(narration_path),
                "slide_timings_path": str(project_dir / "slide_timings.json"),
                "duration_sec": audio_duration,
                "outline": outline,
                "slide_timings": slide_timings,
                "features": ["AI-generated content", "Professional design", "Auto-fetched images", "Charts & diagrams", "Synchronized narration"]
            }
            
            # Add video path if generated
            if video_path:
                response_data["video_path"] = str(video_path)
                response_data["features"].append("MP4 video with synced slides")
            
            return response_data
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Presentation generation failed: {str(e)}\n{traceback.format_exc()}")


# ========== NEW DOCUMENT-BASED ENDPOINTS ==========

from fastapi import UploadFile, File
from .services import document_processor, vector_store


class UploadDocumentRequest(BaseModel):
    """Request model for document upload."""
    project: str


class TopicRequest(BaseModel):
    """Request model for topic-based generation."""
    project: str
    collection_name: str
    topic: str
    lang: str = "en"
    voice: Optional[str] = None
    generate_video: bool = True


@app.post("/upload_document")
async def upload_document(
    file: UploadFile = File(...),
    project: str = "default"
):
    """
    Upload and process PDF/DOCX document.
    
    Extracts text, creates chunks, generates embeddings, stores in ChromaDB.
    Returns topics for user selection.
    """
    try:
        # Validate file type
        if not file.filename:
            raise HTTPException(status_code=400, detail="No filename provided")
        
        suffix = Path(file.filename).suffix.lower()
        if suffix not in ['.pdf', '.docx']:
            raise HTTPException(
                status_code=400,
                detail=f"Unsupported file type: {suffix}. Use PDF or DOCX."
            )
        
        # Create project directory
        project_dir = Path.home() / "Lectures" / project
        project_dir.mkdir(parents=True, exist_ok=True)
        
        # Save uploaded file
        file_path = project_dir / file.filename
        with open(file_path, 'wb') as f:
            content = await file.read()
            f.write(content)
        
        print(f"📁 Saved document: {file_path}")
        
        # Process document
        doc_data = document_processor.process_document(file_path)
        
        # Create collection name from filename (sanitize for ChromaDB)
        # ChromaDB requires: 3-512 chars, [a-zA-Z0-9._-], start/end with [a-zA-Z0-9]
        raw_name = f"{project}_{Path(file.filename).stem}"
        collection_name = re.sub(r'[^a-zA-Z0-9._-]', '_', raw_name)  # Replace invalid chars
        collection_name = re.sub(r'_+', '_', collection_name)  # Collapse multiple underscores
        collection_name = collection_name.strip('._-')  # Remove invalid start/end chars
        
        print(f"📚 Collection name: {collection_name}")
        
        # Get vector store and add documents
        vs = vector_store.get_vector_store(
            persist_directory=str(project_dir / "chroma_db")
        )
        
        vs.add_documents(
            collection_name=collection_name,
            chunks=doc_data['chunks'],
            document_metadata={
                'filename': file.filename,
                'project': project,
                'file_type': doc_data['file_type']
            }
        )
        
        # Save document metadata
        metadata_path = project_dir / f"{Path(file.filename).stem}_metadata.json"
        with open(metadata_path, 'w', encoding='utf-8') as f:
            json.dump({
                'filename': file.filename,
                'collection_name': collection_name,
                'char_count': doc_data['char_count'],
                'chunk_count': doc_data['chunk_count'],
                'topics': doc_data['topics']
            }, f, indent=2)
        
        return {
            "status": "success",
            "filename": file.filename,
            "collection_name": collection_name,
            "char_count": doc_data['char_count'],
            "chunk_count": doc_data['chunk_count'],
            "topics": doc_data['topics'],
            "message": f"Document processed: {doc_data['chunk_count']} chunks indexed"
        }
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"Document upload failed: {str(e)}\n{traceback.format_exc()}"
        )


@app.post("/generate_from_topic")
async def generate_from_topic(request: TopicRequest):
    """
    Generate presentation from document topic using RAG.
    
    1. Search vector store for relevant chunks
    2. Generate presentation from retrieved context
    3. Create slides, audio, video
    """
    try:
        project_dir = Path.home() / "Lectures" / request.project
        project_dir.mkdir(parents=True, exist_ok=True)
        
        # Get vector store
        vs = vector_store.get_vector_store(
            persist_directory=str(project_dir / "chroma_db")
        )
        
        # Search for relevant chunks
        print(f"\n🔍 Searching for topic: {request.topic}")
        search_results = vs.search_similar(
            collection_name=request.collection_name,
            query=request.topic,
            n_results=10  # Get top 10 relevant chunks
        )
        
        if not search_results['chunks']:
            raise HTTPException(
                status_code=404,
                detail=f"No content found for topic: {request.topic}"
            )
        
        # Combine relevant chunks into context
        context = "\n\n".join(search_results['chunks'])
        print(f"✅ Retrieved {len(search_results['chunks'])} relevant chunks")
        print(f"📊 Total context: {len(context)} characters")
        
        # Generate presentation using the context
        # Use the existing presentation generation pipeline but with retrieved context
        request_dict = {
            'project': request.project,
            'topic': f"{request.topic}\n\nContext:\n{context[:10000]}",  # Limit context size
            'lang': request.lang,
            'voice': request.voice,
            'generate_video': request.generate_video
        }
        
        # Call existing presentation generation
        presentation_request = PresentationRequest(**request_dict)
        result = await generate_presentation(presentation_request)
        
        # Add RAG metadata
        result['rag_enabled'] = True
        result['chunks_retrieved'] = len(search_results['chunks'])
        result['relevance_scores'] = [1.0 - d for d in search_results['distances']]
        
        return result
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"Topic-based generation failed: {str(e)}\n{traceback.format_exc()}"
        )


@app.get("/list_collections")
async def list_collections(project: str):
    """List all ChromaDB collections for a project."""
    try:
        project_dir = Path.home() / "Lectures" / project
        chroma_dir = project_dir / "chroma_db"
        
        if not chroma_dir.exists():
            return {"collections": []}
        
        vs = vector_store.get_vector_store(
            persist_directory=str(chroma_dir)
        )
        
        collections = vs.list_collections()
        
        # Get stats for each collection
        collection_stats = []
        for name in collections:
            stats = vs.get_collection_stats(name)
            collection_stats.append(stats)
        
        return {"collections": collection_stats}
        
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"Failed to list collections: {str(e)}"
        )


@app.get("/get_video")
async def get_video(project: str):
    """Get the generated video file path for streaming."""
    try:
        project_dir = Path.home() / "Lectures" / project
        video_path = project_dir / "presentation_video.mp4"
        
        if not video_path.exists():
            raise HTTPException(
                status_code=404,
                detail="Video not found. Generate a presentation first."
            )
        
        return {
            "video_path": str(video_path),
            "exists": True,
            "size_mb": video_path.stat().st_size / (1024 * 1024)
        }
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"Failed to get video: {str(e)}"
        )


@app.get("/get_slide_image")
async def get_slide_image(project: str, slide_index: int):
    """Get a specific slide as an image from the PPTX."""
    try:
        from pptx import Presentation
        from PIL import Image
        import io
        
        project_dir = Path.home() / "Lectures" / project
        pptx_path = project_dir / "presentation.pptx"
        slides_dir = project_dir / "slide_images"
        
        if not pptx_path.exists():
            raise HTTPException(
                status_code=404,
                detail="Presentation not found. Generate a presentation first."
            )
        
        # Create slides directory if it doesn't exist
        slides_dir.mkdir(exist_ok=True)
        
        # Check if image already exists
        image_path = slides_dir / f"slide_{slide_index}.png"
        
        if not image_path.exists():
            # Convert PPTX to images using pptx2pdf approach
            # For now, we'll create placeholder images with slide content
            prs = Presentation(str(pptx_path))
            
            if slide_index < 0 or slide_index >= len(prs.slides):
                raise HTTPException(
                    status_code=404,
                    detail=f"Slide {slide_index} not found. Presentation has {len(prs.slides)} slides."
                )
            
            # Extract slide as image (we'll render text on a canvas)
            slide = prs.slides[slide_index]
            
            # Create image (1920x1080 for 16:9 ratio)
            from PIL import Image, ImageDraw, ImageFont
            img = Image.new('RGB', (1920, 1080), color='#1a1a1a')
            draw = ImageDraw.Draw(img)
            
            # Try to use a nice font, fallback to default
            try:
                title_font = ImageFont.truetype("arial.ttf", 72)
                content_font = ImageFont.truetype("arial.ttf", 48)
            except:
                title_font = ImageFont.load_default()
                content_font = ImageFont.load_default()
            
            y_position = 100
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    
                    # Title vs content
                    if y_position < 200:
                        font = title_font
                        color = '#FFD700'  # Gold for title
                    else:
                        font = content_font
                        color = '#FFFFFF'  # White for content
                    
                    # Word wrap
                    words = text.split()
                    lines = []
                    current_line = []
                    
                    for word in words:
                        current_line.append(word)
                        test_line = ' '.join(current_line)
                        bbox = draw.textbbox((0, 0), test_line, font=font)
                        if bbox[2] - bbox[0] > 1800:  # Max width
                            current_line.pop()
                            if current_line:
                                lines.append(' '.join(current_line))
                            current_line = [word]
                    
                    if current_line:
                        lines.append(' '.join(current_line))
                    
                    # Draw lines
                    for line in lines:
                        bbox = draw.textbbox((0, 0), line, font=font)
                        text_width = bbox[2] - bbox[0]
                        x_position = (1920 - text_width) // 2
                        draw.text((x_position, y_position), line, fill=color, font=font)
                        y_position += 80 if font == title_font else 60
                    
                    y_position += 40
            
            # Save image
            img.save(str(image_path), 'PNG', quality=95)
        
        return {
            "slide_path": str(image_path),
            "exists": True,
            "size_kb": image_path.stat().st_size / 1024
        }
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"Failed to get slide image: {str(e)}\n{traceback.format_exc()}"
        )


@app.get("/get_audio")
async def get_audio(project: str):
    """Get the generated audio file path for interactive lecture."""
    try:
        project_dir = Path.home() / "Lectures" / project
        audio_path = project_dir / "narration.mp3"
        
        if not audio_path.exists():
            raise HTTPException(
                status_code=404,
                detail="Audio not found. Generate a presentation first."
            )
        
        return {
            "audio_path": str(audio_path),
            "exists": True,
            "size_mb": audio_path.stat().st_size / (1024 * 1024)
        }
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"Failed to get audio: {str(e)}"
        )


@app.get("/get_slide_timings")
async def get_slide_timings(project: str):
    """Get the slide timings for accurate slide tracking during video playback."""
    try:
        project_dir = Path.home() / "Lectures" / project
        timings_path = project_dir / "slide_timings.json"
        
        if not timings_path.exists():
            raise HTTPException(
                status_code=404,
                detail="Slide timings not found. Generate a presentation first."
            )
        
        with open(timings_path, "r", encoding="utf-8") as f:
            timings_data = json.load(f)
        
        return timings_data
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"Failed to get slide timings: {str(e)}"
        )


@app.get("/get_animations")
async def get_animations(project: str):
    """Get the interactive animation steps for the lecture player."""
    try:
        project_dir = Path.home() / "Lectures" / project
        animations_path = project_dir / "animations.json"
        
        if not animations_path.exists():
            raise HTTPException(
                status_code=404,
                detail="Animations not found. Generate a presentation first."
            )
        
        with open(animations_path, "r", encoding="utf-8") as f:
            animations_data = json.load(f)
        
        return animations_data
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"Failed to get animations: {str(e)}"
        )


@app.get("/list_projects")
async def list_projects():
    """List all available projects."""
    try:
        lectures_dir = Path.home() / "Lectures"
        if not lectures_dir.exists():
            return {"projects": []}
        
        projects = []
        for item in lectures_dir.iterdir():
            if item.is_dir():
                # Check for metadata files
                video_exists = (item / "presentation_video.mp4").exists()
                pptx_exists = (item / "presentation.pptx").exists()
                
                # Look for document metadata
                metadata_files = list(item.glob("*_metadata.json"))
                collections = []
                for meta_file in metadata_files:
                    with open(meta_file, 'r') as f:
                        meta = json.load(f)
                        collections.append(meta.get('collection_name', ''))
                
                projects.append({
                    'name': item.name,
                    'has_video': video_exists,
                    'has_presentation': pptx_exists,
                    'collections': collections,
                    'path': str(item)
                })
        
        return {"projects": projects}
        
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"Failed to list projects: {str(e)}"
        )


@app.post("/generate_quiz")
async def generate_quiz(request: GenerateQuizRequest):
    """
    Generate interactive quiz for a range of slides.
    
    Args:
        request: Quiz generation request with project, slide range, difficulty
        
    Returns:
        Quiz data with questions, options, answers, and explanations
    """
    from .services.quiz_generator import generate_quiz_for_slides
    
    print(f"\n=== QUIZ GENERATION REQUEST ===")
    print(f"Project: {request.project}")
    print(f"Slide range: {request.slide_start} to {request.slide_end}")
    print(f"Num questions: {request.num_questions}")
    print(f"Difficulty: {request.difficulty}")
    
    try:
        project_path = config.OUTPUT_ROOT / request.project
        
        if not project_path.exists():
            print(f"✗ Project not found: {project_path}")
            raise HTTPException(status_code=404, detail=f"Project not found: {request.project}")
        
        print(f"✓ Project path exists: {project_path}")
        
        quiz_data = generate_quiz_for_slides(
            project_path=project_path,
            slide_range=(request.slide_start, request.slide_end),
            num_questions=request.num_questions,
            difficulty=request.difficulty
        )
        
        print(f"✓ Quiz generated successfully")
        print(f"  Questions generated: {len(quiz_data.get('questions', []))}")
        print(f"  Checkpoint ID: {quiz_data.get('checkpoint_id')}")
        
        # Cache quiz for later retrieval
        quiz_cache_path = project_path / "quiz_cache.json"
        try:
            existing_cache = {}
            if quiz_cache_path.exists():
                with open(quiz_cache_path, 'r', encoding='utf-8') as f:
                    existing_cache = json.load(f)
            
            checkpoint_id = quiz_data["checkpoint_id"]
            existing_cache[checkpoint_id] = quiz_data
            
            with open(quiz_cache_path, 'w', encoding='utf-8') as f:
                json.dump(existing_cache, f, indent=2)
            print(f"✓ Quiz cached to {quiz_cache_path}")
        except Exception as cache_err:
            print(f"⚠ Warning: Failed to cache quiz: {cache_err}")
        
        print(f"=== QUIZ GENERATION COMPLETE ===\n")
        return quiz_data
        
    except FileNotFoundError as e:
        print(f"✗ File not found error: {e}")
        raise HTTPException(status_code=404, detail=str(e))
    except ValueError as e:
        print(f"✗ Value error: {e}")
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        print(f"✗ Quiz generation error: {e}")
        traceback.print_exc()
        raise HTTPException(
            status_code=500,
            detail=f"Quiz generation failed: {str(e)}"
        )


@app.get("/get_quiz_checkpoints")
async def get_quiz_checkpoints(project: str, checkpoint_interval: int = 3):
    """
    Get list of slide indices where quizzes should appear.
    
    Args:
        project: Project name
        checkpoint_interval: Number of slides between quizzes
        
    Returns:
        List of slide indices (0-based) for quiz checkpoints
    """
    from .services.quiz_generator import QuizGenerator
    
    try:
        project_path = config.OUTPUT_ROOT / project
        metadata_path = project_path / "metadata.json"
        
        # If metadata.json doesn't exist, try to create it from script.json
        if not metadata_path.exists():
            script_path = project_path / "script.json"
            if script_path.exists():
                try:
                    with open(script_path, 'r', encoding='utf-8') as f:
                        script_data = json.load(f)
                    
                    # Create metadata from script
                    metadata = {
                        "slides": [
                            {
                                "title": slide.get("title", f"Slide {slide.get('slide_index', i+1)}"),
                                "content": "\n".join(slide.get("points", [])) + "\n" + slide.get("speaker_notes", "")
                            }
                            for i, slide in enumerate(script_data.get("slides", []))
                        ]
                    }
                    
                    # Save metadata
                    with open(metadata_path, 'w', encoding='utf-8') as f:
                        json.dump(metadata, f, indent=2)
                    
                    print(f"Created metadata.json for project: {project}")
                except Exception as e:
                    print(f"Failed to create metadata from script: {e}")
                    raise HTTPException(status_code=404, detail=f"Metadata not found for project: {project}")
            else:
                raise HTTPException(status_code=404, detail=f"Metadata not found for project: {project}")
        
        with open(metadata_path, 'r', encoding='utf-8') as f:
            metadata = json.load(f)
        
        total_slides = len(metadata.get("slides", []))
        
        generator = QuizGenerator()
        checkpoints = generator.generate_quiz_checkpoints(total_slides, checkpoint_interval)
        
        return {
            "project": project,
            "total_slides": total_slides,
            "checkpoint_interval": checkpoint_interval,
            "checkpoints": checkpoints,
            "num_quizzes": len(checkpoints)
        }
        
    except FileNotFoundError as e:
        raise HTTPException(status_code=404, detail=str(e))
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"Failed to get quiz checkpoints: {str(e)}"
        )


@app.post("/submit_quiz_answer")
async def submit_quiz_answer(request: SubmitAnswerRequest):
    """
    Submit and validate a quiz answer, return feedback.
    
    Args:
        request: Answer submission request with project, checkpoint, question, answer
        
    Returns:
        Feedback with correctness, explanation, and correct answer
    """
    try:
        project_path = config.OUTPUT_ROOT / request.project
        quiz_cache_path = project_path / "quiz_cache.json"
        
        if not quiz_cache_path.exists():
            raise HTTPException(status_code=404, detail="Quiz cache not found")
        
        with open(quiz_cache_path, 'r', encoding='utf-8') as f:
            quiz_cache = json.load(f)
        
        if request.checkpoint_id not in quiz_cache:
            raise HTTPException(status_code=404, detail=f"Quiz checkpoint not found: {request.checkpoint_id}")
        
        quiz = quiz_cache[request.checkpoint_id]
        questions = quiz.get("questions", [])
        
        if request.question_index < 0 or request.question_index >= len(questions):
            raise HTTPException(status_code=400, detail="Invalid question index")
        
        question = questions[request.question_index]
        correct_answer = question["correct_answer"]
        is_correct = (request.user_answer.upper() == correct_answer.upper())
        
        # Save user progress
        progress_path = project_path / "study_progress.json"
        try:
            progress = {}
            if progress_path.exists():
                with open(progress_path, 'r', encoding='utf-8') as f:
                    progress = json.load(f)
            
            if request.checkpoint_id not in progress:
                progress[request.checkpoint_id] = {
                    "answers": [],
                    "score": 0,
                    "completed": False
                }
            
            progress[request.checkpoint_id]["answers"].append({
                "question_index": request.question_index,
                "user_answer": request.user_answer,
                "correct_answer": correct_answer,
                "is_correct": is_correct
            })
            
            with open(progress_path, 'w', encoding='utf-8') as f:
                json.dump(progress, f, indent=2)
                
        except Exception as prog_err:
            print(f"Warning: Failed to save progress: {prog_err}")
        
        return {
            "is_correct": is_correct,
            "correct_answer": correct_answer,
            "user_answer": request.user_answer,
            "explanation": question["explanation"],
            "hint": question.get("hint", ""),
            "question_index": request.question_index,
            "checkpoint_id": request.checkpoint_id
        }
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"Failed to submit answer: {str(e)}"
        )


if __name__ == "__main__":
    import uvicorn
    import sys
    
    # Support --port flag for PyInstaller bundle
    port = 8765
    if "--port" in sys.argv:
        idx = sys.argv.index("--port")
        if idx + 1 < len(sys.argv):
            port = int(sys.argv[idx + 1])
    
    uvicorn.run(app, host="127.0.0.1", port=port)
